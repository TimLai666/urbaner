"""
URBANER 商業模式圖 (BMC) + SWOT 分析 PPT 生成
- 含震央選擇 (Epicenter) - 整體 / US / JP 各一份
- 含策略定調卡 (vs 競爭者) - 整體 / US / JP 各一份
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ====== 配色 ======
COLOR_PRIMARY = RGBColor(0x0B, 0x2E, 0x4F)       # 深藍 (URBANER 主色)
COLOR_ACCENT = RGBColor(0xC9, 0x9A, 0x3B)        # 金色 (品味/職人)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF2, 0xEC)      # 米白
COLOR_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
COLOR_LIGHT_GRAY = RGBColor(0xE6, 0xE6, 0xE6)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_RED = RGBColor(0xB8, 0x3A, 0x3A)
COLOR_GREEN = RGBColor(0x3B, 0x7A, 0x4D)
COLOR_BLUE = RGBColor(0x35, 0x6E, 0xA0)
COLOR_AMBER = RGBColor(0xD4, 0x8E, 0x2C)

FONT = 'Microsoft YaHei'

# ====== 投影片尺寸 16:9 ======
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(tf, text, size=14, bold=False, color=COLOR_GRAY, align=PP_ALIGN.LEFT, font=FONT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=12, bold=False,
                color=COLOR_GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    set_text(tf, text, size, bold, color, align, font)
    return tb


def add_multiline_textbox(slide, left, top, width, height, lines, size=11,
                          color=COLOR_GRAY, align=PP_ALIGN.LEFT, font=FONT,
                          fill_color=None, border_color=None):
    """lines: list of (text, bold, size_override) or str"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    if fill_color is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill_color
    if border_color is not None:
        tb.line.color.rgb = border_color
        tb.line.width = Pt(0.75)
    else:
        tb.line.fill.background()

    first = True
    for line in lines:
        if isinstance(line, tuple):
            text, bold, sz = line
        else:
            text, bold, sz = line, False, size

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, border=None, border_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    # remove default shadow
    return shp


def add_header_bar(slide, title, subtitle=None):
    # 頂部主色帶
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.7), COLOR_PRIMARY)
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.5),
                title, size=22, bold=True, color=COLOR_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.3),
                    subtitle, size=12, color=COLOR_ACCENT)
    # 金色細線
    add_rect(slide, Inches(0), Inches(0.7), SLIDE_W, Inches(0.05), COLOR_ACCENT)


def add_footer(slide, page_no, total):
    add_textbox(slide, Inches(11.5), Inches(7.15), Inches(1.6), Inches(0.3),
                f"URBANER × FourSight Lab  |  {page_no}/{total}",
                size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def add_bmc_block(slide, left, top, width, height, title, content_lines,
                  color_bar=COLOR_PRIMARY, fill=COLOR_BG_LIGHT):
    # 主框
    add_rect(slide, left, top, width, height, fill, border=COLOR_LIGHT_GRAY, border_w=1)
    # 上方色條
    add_rect(slide, left, top, width, Inches(0.32), color_bar)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.03), width - Inches(0.16),
                Inches(0.28), title, size=11, bold=True, color=COLOR_WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    # 內容
    body_top = top + Inches(0.36)
    body_height = height - Inches(0.4)
    add_multiline_textbox(slide, left + Inches(0.05), body_top,
                          width - Inches(0.1), body_height,
                          content_lines, size=8.5, color=COLOR_GRAY)


def add_swot_quadrant(slide, left, top, width, height, title, items,
                      bar_color, fill=COLOR_WHITE):
    add_rect(slide, left, top, width, height, fill, border=COLOR_LIGHT_GRAY, border_w=1)
    add_rect(slide, left, top, width, Inches(0.36), bar_color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.04),
                width - Inches(0.2), Inches(0.28),
                title, size=12, bold=True, color=COLOR_WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    add_multiline_textbox(slide, left + Inches(0.08), top + Inches(0.42),
                          width - Inches(0.16), height - Inches(0.5),
                          items, size=9, color=COLOR_GRAY)


# ============================================================
# 開始建構 PPT
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

TOTAL = 18

# ====== Slide 1: 封面 ======
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_PRIMARY)
# 金色裝飾條
add_rect(s, Inches(0.8), Inches(2.2), Inches(0.08), Inches(3.0), COLOR_ACCENT)
add_textbox(s, Inches(1.2), Inches(2.0), Inches(11), Inches(0.6),
            "URBANER 奧本", size=20, bold=True, color=COLOR_ACCENT, font=FONT)
add_textbox(s, Inches(1.2), Inches(2.6), Inches(11), Inches(1.2),
            "商業模式圖 × SWOT 策略分析", size=40, bold=True, color=COLOR_WHITE)
add_textbox(s, Inches(1.2), Inches(3.9), Inches(11), Inches(0.5),
            "整體 × 美國市場 × 日本市場 — 競爭者基準分析",
            size=18, color=COLOR_WHITE)
add_textbox(s, Inches(1.2), Inches(4.6), Inches(11), Inches(0.4),
            "Business Model Canvas (with Epicenter) + SWOT with Strategy Card",
            size=12, color=COLOR_ACCENT)
add_textbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4),
            "Source: URBANER 雙市場 STP 完整報告 (n=11,523 reviews) + 競品評論 + 社群洞察",
            size=10, color=COLOR_LIGHT_GRAY)
add_textbox(s, Inches(1.2), Inches(6.85), Inches(11), Inches(0.4),
            "FourSight Lab × URBANER 產學合作專案  |  資料截止 2026-05-05",
            size=10, color=COLOR_LIGHT_GRAY)


# ====== Slide 2: 目錄 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "目錄 Agenda", "本份報告共 5 個區塊、18 張投影片")
sections = [
    ("01", "競爭者基準與分析方法", "Benchmark scope, S/W 排除規則"),
    ("02", "整體 — 商業模式圖 (含震央)", "Epicenter Selection × 9 elements"),
    ("03", "整體 — SWOT × 策略定調卡", "vs Panasonic / Braun / Wahl / Manscaped"),
    ("04", "美國市場 — BMC × SWOT × 策略", "Gift-Ready × Multi-Function Pro"),
    ("05", "日本市場 — BMC × SWOT × 策略", "精度・耐久・分貝数値化"),
    ("06", "結論與下一步", "三市場戰略地圖總覽"),
]
top = Inches(1.5)
for i, (no, t, sub) in enumerate(sections):
    y = top + Inches(0.85 * i)
    add_rect(s, Inches(0.8), y, Inches(0.7), Inches(0.7), COLOR_ACCENT)
    add_textbox(s, Inches(0.8), y, Inches(0.7), Inches(0.7), no,
                size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(1.7), y + Inches(0.05), Inches(8), Inches(0.35),
                t, size=18, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, Inches(1.7), y + Inches(0.4), Inches(11), Inches(0.3),
                sub, size=11, color=COLOR_GRAY)
add_footer(s, 2, TOTAL)


# ====== Slide 3: 競爭者基準 + 排除規則 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "競爭者基準與 S/W 排除規則",
               "本報告的優勢/劣勢必須是『競爭者沒有』或『URBANER 顯著差異』的項目")

# 競爭者列表 (左)
add_rect(s, Inches(0.4), Inches(1.0), Inches(6.2), Inches(5.8),
         COLOR_BG_LIGHT, border=COLOR_LIGHT_GRAY)
add_textbox(s, Inches(0.55), Inches(1.1), Inches(6), Inches(0.4),
            "競爭者陣容 (依市場)", size=14, bold=True, color=COLOR_PRIMARY)

us_comp = [
    ("美國市場主要競爭者", True, 12),
    ("• Wahl / Andis — 老牌、理髮師取向，正被中國高速無刷馬達品牌侵蝕", False, 9.5),
    ("• Braun (Series 9 Pro) — 高端 Foil Shaver 首選，敏感肌訴求", False, 9.5),
    ("• Manscaped (Lawn Mower 5.0 Ultra) — Body Groomer 標竿、強行銷", False, 9.5),
    ("• Panasonic Arc 5 / Lambdash — Reddit『best bang for buck』", False, 9.5),
    ("• Conair (B008KEJ1LM) — 入門價格錨", False, 9.5),
    ("• ANGFAN / Valano / INVJOY — 中國新進品牌，高 spec 低價", False, 9.5),
    ("", False, 9),
    ("日本市場主要競爭者", True, 12),
    ("• Panasonic — 全品類絕對王者 (ER-GB74-S, ER-GN71, ES-WF41)", False, 9.5),
    ("• Philips Norelco — 海外品牌派、眾多附件", False, 9.5),
    ("• 日立 / Maxell IZUMI — マイベスト#1 鼻毛電剪", False, 9.5),
    ("• Braun — ドイツ信仰、頂級 Foil Shaver", False, 9.5),
    ("• 皇家 (Royal) ET-3 — 手動鼻毛派長期 Bestseller", False, 9.5),
    ("• Pateker (B01K73KIMO) — 寵物電剪日本長期 No.1", False, 9.5),
]
add_multiline_textbox(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(5.2),
                      us_comp, size=10, color=COLOR_GRAY)

# 排除規則 (右)
add_rect(s, Inches(6.8), Inches(1.0), Inches(6.2), Inches(5.8),
         COLOR_WHITE, border=COLOR_LIGHT_GRAY)
add_textbox(s, Inches(6.95), Inches(1.1), Inches(6), Inches(0.4),
            "S/W 排除規則 (來自 SWOT 設計原則)", size=14, bold=True, color=COLOR_PRIMARY)

excl = [
    ("不算 S 的項目 (競爭者也具備)：", True, 11),
    ("• IPX7/IPX8 防水 → 兩市場標配", False, 9.5),
    ("• USB-C 充電 / LED 電量顯示 → 中階以上標配", False, 9.5),
    ("• Multi-attachment (4-7 件) → 全品牌都有", False, 9.5),
    ("• Multi-functional 套組 → 全品牌都有", False, 9.5),
    ("• 鋰電池長續航 → 業界共識", False, 9.5),
    ("", False, 9),
    ("不算 W 的項目 (定位選擇/外部現況)：", True, 11),
    ("• 『價格高於 Conair』→ 定位選擇 (中階定位)", False, 9.5),
    ("• 『預算不足』→ 是中小企業現況、非競爭劣勢", False, 9.5),
    ("• 『中文資源多日文少』→ 已透過 OBM 在地化", False, 9.5),
    ("", False, 9),
    ("S 必須通過四關：", True, 11),
    ("① 市場驗證 (評論數據/銷售支撐)", False, 9.5),
    ("② 相對於競爭者 (非絕對能力)", False, 9.5),
    ("③ 不受公司近期績效污染", False, 9.5),
    ("④ 顧客能感受並影響選擇", False, 9.5),
]
add_multiline_textbox(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(5.2),
                      excl, size=10, color=COLOR_GRAY)
add_footer(s, 3, TOTAL)


# ============================================================
# Section 02: 整體 BMC
# ============================================================

# ====== Slide 4: 整體 — 震央選擇 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "整體 BMC — 震央選擇 (Epicenter Selection)",
               "從四個創新起點中，選定『資源導向』作為整體商業模式設計錨點")

# 四個震央卡片
quads = [
    ("資源導向", "Resource-Driven",
     "47 年 OEM 代工歐美日品牌的製造職人能力 (1977-2014 替 Panasonic 等代工)",
     "✓ 選定：URBANER 整體商模震央",
     "理由：這是兩市場都能用、競爭者複製成本最高的根基",
     COLOR_ACCENT, True),
    ("產品導向", "Product-Driven",
     "以 IPX8 + 0.5mm 精度 + 38 段刻度等規格作為設計原點",
     "× 不選為整體震央",
     "規格 (IPX7/USB-C) 已是全業界標配，無法構成整體差異化",
     COLOR_LIGHT_GRAY, False),
    ("顧客導向", "Customer-Driven",
     "以 US 送禮買家 + JP 精度買家為核心、回推產品",
     "△ 不選為整體震央，但是 US 區域震央",
     "兩市場顧客行為差異過大，不適合作為單一整體錨點",
     COLOR_LIGHT_GRAY, False),
    ("財務導向", "Finance-Driven",
     "OBM 自有品牌毛利 (高於 OEM 代工)、雙引擎收入",
     "× 不選為整體震央",
     "財務是結果而非原因；OBM 轉型已完成 (2020 OBM > OEM)",
     COLOR_LIGHT_GRAY, False),
]
positions = [
    (Inches(0.4), Inches(1.1)),
    (Inches(6.85), Inches(1.1)),
    (Inches(0.4), Inches(4.3)),
    (Inches(6.85), Inches(4.3)),
]
for (label, en, desc, choice, reason, bar, sel), (x, y) in zip(quads, positions):
    fill = COLOR_BG_LIGHT if sel else COLOR_WHITE
    border = COLOR_ACCENT if sel else COLOR_LIGHT_GRAY
    bw = 2.5 if sel else 1
    add_rect(s, x, y, Inches(6.1), Inches(3.05), fill, border=border, border_w=bw)
    add_rect(s, x, y, Inches(6.1), Inches(0.5), bar)
    add_textbox(s, x + Inches(0.15), y + Inches(0.05), Inches(5.9), Inches(0.4),
                f"{label}  ({en})", size=14, bold=True,
                color=COLOR_PRIMARY if sel else COLOR_WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    body = [
        ("核心定義：", True, 10),
        (f"  {desc}", False, 9.5),
        ("", False, 8),
        (choice, True, 11.5),
        (f"  {reason}", False, 9.5),
    ]
    add_multiline_textbox(s, x + Inches(0.15), y + Inches(0.6),
                          Inches(5.8), Inches(2.4), body,
                          size=9.5, color=COLOR_GRAY)
add_footer(s, 4, TOTAL)


# ====== Slide 5: 整體 BMC 9 要素 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "整體商業模式圖 (BMC) — 震央：47 年 OEM 職人製造能力",
               "9 大要素 × 兩市場通用層；US/JP 差異化內容於後續分頁呈現")

# 經典 BMC 九宮格佈局 (橫 5 直 2，但中間是 VP+CR 縱切)
# 我們改用 9 格緊湊版：上層 5 格 (KP|KA|VP|CR|CS) 下層 2 格 (CH+REV) + 中段 KR
# 為簡化，使用 6 列 × 2 row 結構 + 底部 cost+revenue
margin_x = Inches(0.3)
margin_y = Inches(1.05)
gap = Inches(0.08)

# 標準 BMC 五縱欄佈局
top_h = Inches(2.85)
mid_h = Inches(2.85)
bot_h = Inches(1.35)

col1_w = Inches(2.55)  # KP
col2_w = Inches(2.55)  # KA / KR
col3_w = Inches(2.6)   # VP
col4_w = Inches(2.55)  # CR / CH
col5_w = Inches(2.55)  # CS

# 計算 X 座標
x1 = margin_x
x2 = x1 + col1_w + gap
x3 = x2 + col2_w + gap
x4 = x3 + col3_w + gap
x5 = x4 + col4_w + gap

# Top row 高度 = (top_h-gap)/2
half_h = Inches(1.39)

# 1. KP (Key Partners)
add_bmc_block(s, x1, margin_y, col1_w, top_h, "① 關鍵合作夥伴 KP", [
    ("• 代工夥伴：原 47 年合作的歐美日品牌 (Panasonic 等)", False, 8),
    ("  ─ 持續代工保留現金流", False, 8),
    ("• Amazon (US/JP) — 主力銷售平台", False, 8),
    ("• 樂天市場 — 台灣廠商稀缺直營資格", False, 8),
    ("• 群眾募資平台 (日本) — 品牌曝光", False, 8),
    ("• 中國/越南 OEM 製造夥伴 — 補產能", False, 8),
    ("• 日本物流商 — 在地倉儲與保固", False, 8),
    ("• 跨境金物流支付服務商", False, 8),
])

# 2. KA (Key Activities)
add_bmc_block(s, x2, margin_y, col2_w, half_h, "② 關鍵活動 KA", [
    ("• 產品設計 + 模具開發 (留台)", False, 8),
    ("• Amazon Listing × A+ 內容 × 廣告操作", False, 8),
    ("• 雙語 (英/日) 在地化內容生產", False, 8),
    ("• 評論監控 + 屬性回饋至 R&D", False, 8),
    ("• OEM × OBM 雙引擎平衡管理", False, 8),
])

# KR (Key Resources) - 中下
add_bmc_block(s, x2, margin_y + half_h + gap, col2_w, half_h, "③ 關鍵資源 KR", [
    ("• 47 年 OEM 製造 know-how × 模具", False, 8),
    ("• 11,523 則自有 Amazon 評論資產", False, 8),
    ("• 樂天直營資格 (台廠稀缺)", False, 8),
    ("• 林雷鈞家族二代接班團隊", False, 8),
    ("• 日本 URBANER 跨境官網與品牌", False, 8),
])

# 4. VP (Value Proposition) — 中央，跨上下
add_bmc_block(s, x3, margin_y, col3_w, top_h, "④ 價值主張 VP", [
    ("『品味源自細節』", True, 11),
    ("Bright Idea, Great Value", True, 9),
    ("", False, 6),
    ("• 日系工藝 × 台灣設計 × 親民售價", False, 8.5),
    ("  (Conair 入門以上，Braun 旗艦以下)", False, 8.5),
    ("", False, 5),
    ("• 一機多用 + 完整附件 (4-7 件)", False, 8.5),
    ("• IPX7 全水洗、敏感肌友善", False, 8.5),
    ("• 47 年代工歐美日大廠的職人精神", False, 8.5),
    ("• Beard / Nose-Ear 雙旗艦明確", False, 8.5),
    ("• 雙市場差異化主打 SKU", False, 8.5),
    ("  US: B0FL267TCG (Beard)", False, 8.5),
    ("  JP: B0GBWZBMS5 (Nose/Ear)", False, 8.5),
], color_bar=COLOR_ACCENT, fill=COLOR_BG_LIGHT)

# 5. CR (Customer Relationships)
add_bmc_block(s, x4, margin_y, col4_w, half_h, "⑤ 顧客關係 CR", [
    ("• Amazon 評論回應 + 改版迭代", False, 8),
    ("• 樂天直營店 1 年保固承諾", False, 8),
    ("• 日本『のし対応 + 禮盒包裝』在地服務", False, 8),
    ("• 跨境官網 FAQ (日文/英文)", False, 8),
    ("• 募資平台 backer 早鳥社群", False, 8),
])

# CH (Channels)
add_bmc_block(s, x4, margin_y + half_h + gap, col4_w, half_h, "⑥ 通路 CH", [
    ("• Amazon US：amazon.com/urbaner", False, 8),
    ("• Amazon JP：amazon.co.jp/URBANER", False, 8),
    ("• 樂天市場 (JP)：台廠直營", False, 8),
    ("• 跨境官網：urbaner.jp/ja", False, 8),
    ("• 群眾募資 (JP) + 台灣本地電商/實體", False, 8),
])

# 7. CS (Customer Segments)
add_bmc_block(s, x5, margin_y, col5_w, top_h, "⑦ 目標客群 CS", [
    ("整體目標：", True, 9),
    ("中階以上注重品味的男士/家庭顧客", False, 8.5),
    ("", False, 5),
    ("US 主力：", True, 9),
    ("• 父親節/聖誕節送禮買家 (Q2/Q4)", False, 8),
    ("• 76.5% Beard 多功能套組買家", False, 8),
    ("", False, 5),
    ("JP 主力：", True, 9),
    ("• 91.6% 鼻毛+鬍鬚 CP 值 / 易用", False, 8),
    ("• 7.8% 高精度刻度玩家", False, 8),
    ("", False, 5),
    ("跨市場：", True, 9),
    ("• 寵物/嬰兒類補強客群", False, 8),
])

# 8. CST (Cost Structure) - 底部左半
add_bmc_block(s, x1, margin_y + top_h + gap, Inches(6.31), bot_h,
              "⑧ 成本結構 Cost Structure", [
    ("• 製造 (台/中越 OEM 廠)：原料 + 人工 + 模具攤提", False, 8),
    ("• Amazon 平台費 + 廣告費 (PPC/DSP) — 占銷售 15-25%", False, 8),
    ("• 跨境物流 + 海關 + 倉儲 (FBA / 樂天倉)", False, 8),
    ("• R&D + 雙語在地化 + 募資代操費用", False, 8),
])

# 9. REV (Revenue Streams) - 底部右半
add_bmc_block(s, x3 + Inches(0.05), margin_y + top_h + gap, Inches(6.31), bot_h,
              "⑨ 收益流 Revenue Streams", [
    ("• 雙引擎：OBM (URBANER) 品牌銷售 (主) + OEM 代工 (現金流)", False, 8),
    ("• 美國 Amazon 銷售 (USD)、日本 Amazon + 樂天 (JPY)", False, 8),
    ("• 替換刀頭/耗材後續銷售 (recurring)", False, 8),
    ("• 募資專案 (高曝光 + 預售現金)", False, 8),
])
add_footer(s, 5, TOTAL)


# ============================================================
# Section 03: 整體 SWOT
# ============================================================

# ====== Slide 6: 整體 SWOT 矩陣 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "整體 SWOT — 與 Panasonic / Braun / Wahl / Manscaped 對比",
               "已剃除『競爭者也有』的項目 (IPX7/USB-C/多附件等共通標配不列入 S/W)")

S_items = [
    ("S1  47 年 OEM 製造職人能力", True, 10),
    ("    替 Panasonic 等代工累積，新進品牌(ANGFAN, Valano)無此底蘊", False, 8.5),
    ("S2  日本樂天直營店資格", True, 10),
    ("    台灣廠商稀缺資源；多數 Wahl/Andis/Conair 沒有", False, 8.5),
    ("S3  US/JP 雙市場原生 Listing 雙語經營", True, 10),
    ("    Manscaped 主力 US；皇家/Maxell IZUMI 主力 JP；URBANER 雙線並進", False, 8.5),
    ("S4  Beard + Nose/Ear 雙旗艦 Hero SKU 已成形", True, 10),
    ("    11,523 則自有評論累積為 R&D 資產", False, 8.5),
    ("S5  OBM + OEM 雙引擎商模降低風險", True, 10),
    ("    純品牌如 Manscaped 無 OEM 緩衝", False, 8.5),
]
W_items = [
    ("W1  品牌知名度顯著低於 Panasonic / Braun / Wahl", True, 10),
    ("    JP 知恵袋/価格.com 幾乎無 URBANER 討論", False, 8.5),
    ("W2  缺乏專家/權威背書", True, 10),
    ("    Panasonic『美容クリニック推奨』、ER-WF41 醫療指定 — URBANER 無", False, 8.5),
    ("W3  Foil Shaver / Pet / Body Groomer 品類薄弱", True, 10),
    ("    Beard/Nose 過度集中；Braun-9、Manscaped 5.0、Casfuy 各有壓倒地位", False, 8.5),
    ("W4  社群/Reddit/マイベスト 排名缺席", True, 10),
    ("    競品多在 Reddit/Kakaku.com 高頻被推薦", False, 8.5),
    ("W5  創新節奏較慢", True, 10),
    ("    中國新進(ANGFAN, Valano)用『高速無刷馬達 + LCD』快速跑出市佔", False, 8.5),
]
O_items = [
    ("O1  『靜音 (低分貝)』全品類大趨勢 (寵物/嬰兒/男士)", True, 10),
    ("    URBANER 可用 dB 數字化 + 馬達技術切入", False, 8.5),
    ("O2  美國 Q2(父親節) + Q4(聖誕) 送禮季", True, 10),
    ("    gift_suitability_men ANOVA F=1299 為最強區隔屬性", False, 8.5),
    ("O3  日本父の日 + のし禮品文化", True, 10),
    ("    樂天直營有資格做『正規品 + 1年保証』訴求", False, 8.5),
    ("O4  消費者搜尋『Manscaped 替代品』意圖高", True, 10),
    ("    Reddit 用戶主動找平價代替品", False, 8.5),
    ("O5  跨境電商基礎設施成熟 (Amazon Global, FBA)", True, 10),
    ("    台灣中小企業可低成本進入兩市場", False, 8.5),
]
T_items = [
    ("T1  中國新進品牌『高 spec/低價』夾擊", True, 10),
    ("    ANGFAN, Valano, INVJOY 在 US 已侵蝕 Wahl 市佔", False, 8.5),
    ("T2  Panasonic / Braun 行銷預算碾壓", True, 10),
    ("    JP 家電量販店通路 + 廣告投放遠超 URBANER", False, 8.5),
    ("T3  Amazon 平台規則 + 廣告費上漲", True, 10),
    ("    PPC CPC 持續上升、A9 演算法變動風險", False, 8.5),
    ("T4  匯率波動 (USD/JPY)", True, 10),
    ("    日圓貶值有利日本銷售但壓縮人民幣/台幣換匯利潤", False, 8.5),
    ("T5  山寨 / 跨境平行輸入侵蝕", True, 10),
    ("    Amazon 假貨檢舉週期長", False, 8.5),
]

quad_w = Inches(6.4)
quad_h = Inches(2.85)
y_top = Inches(1.05)
y_bot = y_top + quad_h + Inches(0.1)

add_swot_quadrant(s, Inches(0.3), y_top, quad_w, quad_h,
                  "S — 優勢 (vs 上述競爭者)", S_items, COLOR_GREEN)
add_swot_quadrant(s, Inches(6.83), y_top, quad_w, quad_h,
                  "W — 劣勢 (vs 上述競爭者)", W_items, COLOR_RED)
add_swot_quadrant(s, Inches(0.3), y_bot, quad_w, quad_h,
                  "O — 機會 (看得到 ✓ 吃得到 ✓ 對手難以同樣掌握)", O_items, COLOR_BLUE)
add_swot_quadrant(s, Inches(6.83), y_bot, quad_w, quad_h,
                  "T — 威脅 (立即且迫切)", T_items, COLOR_AMBER)
add_footer(s, 6, TOTAL)


# ====== Slide 7: 整體 SWOT 策略碰撞 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "整體 SWOT 策略碰撞 — TOWS 矩陣 (每條策略需標明 Sx×Oy 來源)",
               "整體態勢：S>W (有獨特資源) × O>T (機會大於威脅) → 主象限：SO 乘勝追擊")

so_items = [
    ("[S1×O3] 以『47年代工歷練的日系工藝』結合 樂天直営正規品 +", True, 9.5),
    ("        1年保証 + のし対応，搶 JP 父の日禮品檔期", False, 8.5),
    ("[S4×O2] 集中 Hero SKU B0FL267TCG，做 Father's Day 禮盒版", True, 9.5),
    ("        (含禮品袋 + 賀卡) 直攻 US Q2 送禮潮", False, 8.5),
    ("[S2×O3] 樂天直營資格 + 1年保証 + 跨境官網結合，做 JP 信任", True, 9.5),
    ("        雙保險 (對抗 Panasonic 通路深度)", False, 8.5),
    ("[S3×O5] 雙市場 Listing 同步迭代、評論回饋環跑得比競品快", True, 9.5),
]
st_items = [
    ("[S5×T1] OEM 引擎吸納成本壓力、OBM 維持品質與毛利，", True, 9.5),
    ("        對中國低價競品做差異化區隔 (Quality > Price)", False, 8.5),
    ("[S1×T2] 把『47 年代工歐美日品牌』做成品牌故事影片，", True, 9.5),
    ("        對抗 Panasonic 廣告聲量 (用故事 vs 預算)", False, 8.5),
    ("[S4×T3] 集中廣告預算於 Hero SKU 而非廣撒，提升 ROAS", True, 9.5),
    ("        以對抗 Amazon CPC 上升", False, 8.5),
]
wo_items = [
    ("[W2×O1] 與獸醫/理髮師/美容診所合作做『専門家監修』背書，", True, 9.5),
    ("        切入靜音 (dB 數值化) 寵物/嬰兒新品類", False, 8.5),
    ("[W4×O4] 主動投放 Reddit / マイベスト / 価格.com 評測，", True, 9.5),
    ("        卡『Manscaped 平價替代』『Wahl 替代品』搜尋意圖", False, 8.5),
    ("[W1×O5] 用 Amazon Choice / Best Seller 徽章 + 評論積累", True, 9.5),
    ("        逐步建立平台原生品牌力", False, 8.5),
]
wt_items = [
    ("[W3×T1] 暫緩 Foil Shaver / Body Groomer 自有品牌投入，", True, 9.5),
    ("        改以 OEM 代工形式維持產線稼動率", False, 8.5),
    ("[W5×T1] 縮短新品 NPI 週期 (從 18 個月→ 12 個月)，", True, 9.5),
    ("        每年至少 2 款新 Hero SKU 對抗中國新進", False, 8.5),
    ("[W4×T5] 強化『URBANER 直営正規品』水印 + 序號系統，", True, 9.5),
    ("        對抗山寨平行輸入", False, 8.5),
]

quad_w2 = Inches(6.4)
quad_h2 = Inches(2.85)
add_swot_quadrant(s, Inches(0.3), Inches(1.05), quad_w2, quad_h2,
                  "SO 乘勝追擊 (主象限) — 攻擊/維持領先", so_items, COLOR_GREEN)
add_swot_quadrant(s, Inches(6.83), Inches(1.05), quad_w2, quad_h2,
                  "ST 守株待兔 — 防禦/移轉威脅", st_items, COLOR_AMBER)
add_swot_quadrant(s, Inches(0.3), Inches(4.0), quad_w2, quad_h2,
                  "WO 策略聯盟 — 補強弱勢、抓住機會", wo_items, COLOR_BLUE)
add_swot_quadrant(s, Inches(6.83), Inches(4.0), quad_w2, quad_h2,
                  "WT 置之死地後生 — 退守/避險", wt_items, COLOR_RED)
add_footer(s, 7, TOTAL)


# ====== Slide 8: 整體 — 策略定調卡 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "整體 — 策略定調卡 (Strategy Card)",
               "從 13 條候選策略中選定 1 條主策略 + 支援/暫緩/放棄清單")

# 左欄：候選策略總覽 + 評分
add_rect(s, Inches(0.3), Inches(1.0), Inches(6.4), Inches(6.0),
         COLOR_BG_LIGHT, border=COLOR_LIGHT_GRAY)
add_textbox(s, Inches(0.45), Inches(1.1), Inches(6.2), Inches(0.4),
            "候選策略總覽 (主象限 SO 評分)", size=14, bold=True, color=COLOR_PRIMARY)

candidates = [
    ("[S1×O3]  日系工藝 × 樂天直営 × 父の日禮品", True, 10.5),
    ("    可行性 3 / 影響力 3 / 時效性 3 = 9 分 ★主策略候選", False, 9),
    ("", False, 5),
    ("[S4×O2]  Hero SKU 美國父親節禮盒版", True, 10.5),
    ("    可行性 3 / 影響力 2 / 時效性 3 = 8 分", False, 9),
    ("", False, 5),
    ("[S2×O3]  樂天 1年保証 × 跨境官網信任雙保險", True, 10.5),
    ("    可行性 3 / 影響力 2 / 時效性 2 = 7 分", False, 9),
    ("", False, 5),
    ("[S3×O5]  雙市場評論回饋環", True, 10.5),
    ("    可行性 2 / 影響力 2 / 時效性 2 = 6 分 (基礎建設)", False, 9),
    ("", False, 5),
    ("評分維度說明", True, 10),
    ("• 可行性：現有資源能否啟動", False, 8.5),
    ("• 影響力：對競爭地位改變幅度", False, 8.5),
    ("• 時效性：機會/威脅窗口是否仍開", False, 8.5),
]
add_multiline_textbox(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.4),
                      candidates, size=10, color=COLOR_GRAY)

# 右欄：主策略 + 支援/暫緩/放棄
add_rect(s, Inches(6.85), Inches(1.0), Inches(6.2), Inches(6.0),
         COLOR_WHITE, border=COLOR_LIGHT_GRAY)
add_rect(s, Inches(6.85), Inches(1.0), Inches(6.2), Inches(0.4), COLOR_ACCENT)
add_textbox(s, Inches(7.0), Inches(1.04), Inches(6), Inches(0.32),
            "✓ 主策略 (整體最終選定)", size=13, bold=True, color=COLOR_WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

main_strat = [
    ("來源：[S1×O3]  日系職人工藝 × 樂天 × 父の日禮品", True, 11),
    ("", False, 4),
    ("策略描述：", True, 10),
    ("以『47 年代工歐美日大廠』為品牌故事核心，搭配", False, 9.5),
    ("樂天直營正規品 + 1 年保証 + のし禮盒，主攻日本", False, 9.5),
    ("父の日 (6 月) 與年末ギフト檔期；同步用『日系工藝』", False, 9.5),
    ("故事複用於美國 Father's Day Q2 與 Q4 聖誕。", False, 9.5),
    ("", False, 5),
    ("為何優於其他：", True, 10),
    ("• 同時槓桿 S1 (代工 know-how) + S2 (樂天) + S3 (雙語)", False, 9.5),
    ("• 直攻 ANOVA Top 1 區隔屬性 gift_suitability_men", False, 9.5),
    ("• 競爭者 Wahl/Manscaped 無 JP 樂天通路、無職人故事可說", False, 9.5),
    ("", False, 5),
    ("🔁 支援策略：[W4×O4] Reddit/マイベスト 評測投放", True, 9.5),
    ("⏸ 暫緩策略：[W3×T1] Foil/Body Groomer 自有品牌 — 等待 ", True, 9.5),
    ("    ：Beard/Nose 雙旗艦每年營收 > USD 5M 後再啟動", False, 8.5),
    ("❌ 放棄策略：[W2×O1] 自建獸醫/醫療診所聯盟 — 原因：", True, 9.5),
    ("    投資回收期 > 24 個月、現金流壓力過大", False, 8.5),
]
add_multiline_textbox(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(5.4),
                      main_strat, size=9.5, color=COLOR_GRAY)
add_footer(s, 8, TOTAL)


# ============================================================
# Section 04: 美國市場 BMC + SWOT
# ============================================================

# ====== Slide 9: 美國 — 震央 + BMC ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "美國市場 — BMC 震央：顧客導向 (Customer-Driven)",
               "ANOVA Top 1 屬性 gift_suitability_men (F=1299) → 從『送禮買家』回推產品/通路/收益")

# 震央說明條
add_rect(s, Inches(0.3), Inches(1.0), Inches(12.7), Inches(0.6),
         COLOR_ACCENT)
add_textbox(s, Inches(0.45), Inches(1.05), Inches(12.5), Inches(0.5),
            "震央：顧客導向 — 美國 76.5% 區隔 + 19.5% 區隔皆強烈受『送禮場景』驅動，"
            "從顧客 → VP → 產品/通路/收益鏈一氣呵成",
            size=11, bold=True, color=COLOR_WHITE, anchor=MSO_ANCHOR.MIDDLE)

# BMC 9 元素 (簡化排版)
margin_y2 = Inches(1.75)
top_h2 = Inches(2.8)
bot_h2 = Inches(1.3)
half_h2 = Inches(1.36)

add_bmc_block(s, x1, margin_y2, col1_w, top_h2, "① KP", [
    ("• Amazon US (FBA + PPC + DSP)", False, 8.5),
    ("• 美國節慶禮品包裝供應商", False, 8.5),
    ("• 美國當地 Influencer / YouTuber", False, 8.5),
    ("  (Beard care 領域)", False, 8.5),
    ("• Manscaped 替代品評測媒體", False, 8.5),
])
add_bmc_block(s, x2, margin_y2, col2_w, half_h2, "② KA", [
    ("• 父親節/聖誕節禮盒企劃", False, 8),
    ("• Listing A+ 內容 (Gift Ready 主訴求)", False, 8),
    ("• PPC 預算季節性集中投放 Q2/Q4", False, 8),
])
add_bmc_block(s, x2, margin_y2 + half_h2 + gap, col2_w, half_h2, "③ KR", [
    ("• Hero SKU B0FL267TCG (Beard)", False, 8),
    ("• 4,360 則 US 自有評論", False, 8),
    ("• 多功能 5-in-1 / 7-in-1 套組", False, 8),
])
add_bmc_block(s, x3, margin_y2, col3_w, top_h2, "④ VP — 美國版", [
    ("『The Perfect Gift for Him』", True, 11),
    ("一台搞定多功能 × 直接送禮", True, 9),
    ("", False, 5),
    ("• 7-in-1 多功能套組 (主圖右上徽章)", False, 8.5),
    ("• 父親節/聖誕節 ready 禮盒包裝", False, 8.5),
    ("• 日系 OEM 工藝 (vs Wahl/ANGFAN)", False, 8.5),
    ("• IPX7 + Wet/Dry + LCD 電量", False, 8.5),
    ("• 介於 Conair (入門) 與 Braun 9", False, 8.5),
    ("  Pro 旗艦 ($300+) 之間的中階", False, 8.5),
    ("  $60-120 甜蜜帶", False, 8.5),
    ("• 鎖定 Segment 1+2 (96% 顧客)", False, 8.5),
], color_bar=COLOR_ACCENT, fill=COLOR_BG_LIGHT)
add_bmc_block(s, x4, margin_y2, col4_w, half_h2, "⑤ CR", [
    ("• 評論回應 (Beard 與 Ear Hair 痛點)", False, 8),
    ("• 退換貨政策對齊 Amazon 預期", False, 8),
    ("• Father's Day Email 觸發行銷", False, 8),
])
add_bmc_block(s, x4, margin_y2 + half_h2 + gap, col4_w, half_h2, "⑥ CH", [
    ("• Amazon US (主) — amazon.com/urbaner", False, 8),
    ("• Amazon Live + Sponsored Video", False, 8),
    ("• Reddit/YouTube Beard 圈評測投放", False, 8),
])
add_bmc_block(s, x5, margin_y2, col5_w, top_h2, "⑦ CS — 美國", [
    ("Segment 2 (76.5%, n=3,337)", True, 9),
    ("Ear Hair Trimming Function 主軸", False, 8),
    ("avg★ = 3.18 → 痛點救火優先", False, 8),
    ("", False, 4),
    ("Segment 1 (19.5%, n=850)", True, 9),
    ("Gift Suitability Men 主軸", False, 8),
    ("avg★ = 3.68 → 送禮場景擴張", False, 8),
    ("", False, 4),
    ("Segment 3 (4.0%, n=173)", True, 9),
    ("Power Source / 高分群 (★4.48)", False, 8),
    ("→ 重點維護的口碑核心", False, 8),
])
add_bmc_block(s, x1, margin_y2 + top_h2 + gap, Inches(6.31), bot_h2,
              "⑧ 成本結構", [
    ("• Amazon FBA 物流 + 平台費 (~15%) + PPC (~10-15%)", False, 8),
    ("• 美國節慶包裝設計與打樣費用 (Q1/Q3 提前)", False, 8),
    ("• Influencer / Reddit/YouTube 評測費用", False, 8),
])
add_bmc_block(s, x3 + Inches(0.05), margin_y2 + top_h2 + gap, Inches(6.31), bot_h2,
              "⑨ 收益流", [
    ("• 主：Beard / Mustache Trimmer 系列銷售 (Hero B0FL267TCG)", False, 8),
    ("• Q2 父親節 + Q4 聖誕禮盒檔期 = 全年營收 50-60%", False, 8),
    ("• 替換刀頭 / 油 / 配件後續銷售 (高毛利 recurring)", False, 8),
])
add_footer(s, 9, TOTAL)


# ====== Slide 10: 美國 SWOT ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "美國 SWOT — 與 Wahl / Andis / Manscaped / Braun / Conair / ANGFAN 對比",
               "排除『美國競爭者也有』的項目；只列 URBANER 在美國市場相對差異化")

US_S = [
    ("S1  日系 OEM 工藝故事 (47 年代工)", True, 10),
    ("    Wahl/Andis/ANGFAN/Manscaped 都無此底蘊", False, 8.5),
    ("S2  Beard Hero SKU B0FL267TCG 已成形", True, 10),
    ("    距理想點 1.601；3,337 則 Beard 評論底量", False, 8.5),
    ("S3  中階價格定位 ($60-120) — 甜蜜帶", True, 10),
    ("    Conair 入門以上、Braun 9 Pro ($300) 以下", False, 8.5),
    ("S4  雙語跨境營運 (US+JP 兩市場現金流分散風險)", True, 10),
    ("    Manscaped 主力 US；Wahl 主力 US；URBANER 雙線", False, 8.5),
]
US_W = [
    ("W1  Reddit / YouTube / B&B 論壇能見度低", True, 10),
    ("    競品 Wahl/Andis/Manscaped 在 Reddit 高頻被討論", False, 8.5),
    ("W2  Body Groomer 品類缺席", True, 10),
    ("    Manscaped Lawn Mower 5.0 Ultra『無 guard 也不傷』獨佔", False, 8.5),
    ("W3  Foil Shaver 品類薄弱", True, 10),
    ("    Braun Series 9 + Panasonic Arc 5 雙佔", False, 8.5),
    ("W4  Pet 品類在美無 Hero SKU", True, 10),
    ("    Casfuy 終身保固已成『最不會被閒置』口碑", False, 8.5),
    ("W5  缺 Q2/Q4 節慶以外的買點故事", True, 10),
    ("    淡季 (Q1/Q3) 銷售落差過大", False, 8.5),
]
US_O = [
    ("O1  父親節 (6/15) + 聖誕 (12/25) 送禮潮", True, 10),
    ("    gift_suitability_men ANOVA F=1299 (Top 1 區隔)", False, 8.5),
    ("O2  Wahl 漸被 ANGFAN 等高速無刷馬達淘汰", True, 10),
    ("    用戶主動轉移 → URBANER 可乘機卡位", False, 8.5),
    ("O3  Manscaped 替代品搜尋量大", True, 10),
    ("    Reddit『Manscaped alternatives』高頻搜尋", False, 8.5),
    ("O4  Beard 多功能套組需求 (num_grooming_functions F=1193)", True, 10),
    ("    7-in-1 套組形態為市場期待", False, 8.5),
]
US_T = [
    ("T1  中國新進 (ANGFAN/Valano/INVJOY) 高 spec/低價", True, 10),
    ("    LCD + USB-C + 高速馬達同等規格價格更低", False, 8.5),
    ("T2  Manscaped 行銷預算碾壓 (Super Bowl 廣告)", True, 10),
    ("    Body Groomer 心智占有極高", False, 8.5),
    ("T3  Amazon US PPC CPC 持續上升", True, 10),
    ("    Beard 類目 ACoS 已到 30%+", False, 8.5),
    ("T4  美中關稅政策不確定性", True, 10),
    ("    若關稅升至 25%+ 將直接侵蝕毛利", False, 8.5),
]

add_swot_quadrant(s, Inches(0.3), Inches(1.05), quad_w, quad_h,
                  "S — 美國優勢 (vs Wahl/Manscaped/Braun)", US_S, COLOR_GREEN)
add_swot_quadrant(s, Inches(6.83), Inches(1.05), quad_w, quad_h,
                  "W — 美國劣勢", US_W, COLOR_RED)
add_swot_quadrant(s, Inches(0.3), Inches(4.0), quad_w, quad_h,
                  "O — 美國機會", US_O, COLOR_BLUE)
add_swot_quadrant(s, Inches(6.83), Inches(4.0), quad_w, quad_h,
                  "T — 美國威脅", US_T, COLOR_AMBER)
add_footer(s, 10, TOTAL)


# ====== Slide 11: 美國 SWOT 策略碰撞 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "美國 SWOT 策略碰撞 — TOWS 矩陣",
               "態勢：S>W (有日系工藝故事) × O>T (送禮潮 + Wahl 退場機會) → 主象限：SO")

us_so = [
    ("[S1×O1] 把『47年代工日系工藝』包裝成 Father's Day 主題影片", True, 9.5),
    ("        ＋禮盒，直攻 Q2/Q4 送禮潮 ★主策略候選", False, 8.5),
    ("[S2×O4] B0FL267TCG 升級成 7-in-1 套組版主推 Q4", True, 9.5),
    ("[S3×O2] Wahl 退場帶動的價格甜蜜帶 ($60-120) 卡位", True, 9.5),
    ("[S1×O3] 製作『Manscaped 替代』比較圖、強調日系工藝", True, 9.5),
]
us_st = [
    ("[S3×T1] 中階定位避開中國 < $40 紅海，", True, 9.5),
    ("        強化『日系工藝 + 1年保固』溢價合理化", False, 8.5),
    ("[S4×T4] 美中關稅升高時、加速移轉至 OEM 越南/台灣產線", True, 9.5),
    ("        對抗供應鏈成本衝擊 (利用 OBM+OEM 雙引擎彈性)", False, 8.5),
    ("[S2×T3] PPC 集中於 Hero SKU 而非廣撒，提升 ROAS", True, 9.5),
]
us_wo = [
    ("[W1×O3] 主動經營 Reddit/r/wicked_edge/r/beards/r/Manscaping", True, 9.5),
    ("        + B&B 論壇，建立『Manscaped 平價替代』口碑", False, 8.5),
    ("[W5×O1] 開發母親節 (5月) 女士系列 + 寵物父親節迷你版", True, 9.5),
    ("        填補淡季缺口", False, 8.5),
]
us_wt = [
    ("[W2×T2] Body Groomer 品類『不投入自有品牌』", True, 9.5),
    ("        改以 OEM 接 Manscaped 競品的單，避開正面對決", False, 8.5),
    ("[W4×T1] Pet 品類以 OEM 為主、暫停 OBM 投入", True, 9.5),
    ("        避免在 Casfuy 主場硬碰", False, 8.5),
]

add_swot_quadrant(s, Inches(0.3), Inches(1.05), quad_w2, quad_h2,
                  "SO 美國主象限 — 乘勝追擊", us_so, COLOR_GREEN)
add_swot_quadrant(s, Inches(6.83), Inches(1.05), quad_w2, quad_h2,
                  "ST — 守株待兔", us_st, COLOR_AMBER)
add_swot_quadrant(s, Inches(0.3), Inches(4.0), quad_w2, quad_h2,
                  "WO — 策略聯盟", us_wo, COLOR_BLUE)
add_swot_quadrant(s, Inches(6.83), Inches(4.0), quad_w2, quad_h2,
                  "WT — 退守/避險", us_wt, COLOR_RED)
add_footer(s, 11, TOTAL)


# ====== Slide 12: 美國策略定調卡 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "美國 — 策略定調卡",
               "整體態勢 S>W × O>T → SO 主象限 → 主策略：Gift-Ready × Multi-Function Pro")

# 候選評分
add_rect(s, Inches(0.3), Inches(1.0), Inches(6.4), Inches(6.0),
         COLOR_BG_LIGHT, border=COLOR_LIGHT_GRAY)
add_textbox(s, Inches(0.45), Inches(1.1), Inches(6.2), Inches(0.4),
            "美國候選策略評分 (主象限 SO)", size=14, bold=True, color=COLOR_PRIMARY)

us_cand = [
    ("[S1×O1]  Father's Day 日系工藝禮盒", True, 10.5),
    ("    可行性 3 / 影響力 3 / 時效性 3 = 9 ★ 最高分", False, 9),
    ("", False, 5),
    ("[S2×O4]  Hero SKU 升級 7-in-1 套組", True, 10.5),
    ("    可行性 3 / 影響力 3 / 時效性 2 = 8", False, 9),
    ("", False, 5),
    ("[S3×O2]  價格甜蜜帶卡位 (Wahl 退場區間)", True, 10.5),
    ("    可行性 3 / 影響力 2 / 時效性 3 = 8", False, 9),
    ("", False, 5),
    ("[S1×O3]  Manscaped 替代比較行銷", True, 10.5),
    ("    可行性 2 / 影響力 2 / 時效性 3 = 7", False, 9),
    ("", False, 5),
    ("關鍵節點時程", True, 10.5),
    ("• 1-2 月：禮盒打樣與供應商談判", False, 9),
    ("• 3-4 月：Listing 改版 + A+ 內容更新", False, 9),
    ("• 5-6 月：Father's Day 廣告主檔期", False, 9),
    ("• 9-11 月：複用素材跑黑五/聖誕", False, 9),
]
add_multiline_textbox(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.4),
                      us_cand, size=10, color=COLOR_GRAY)

# 主策略
add_rect(s, Inches(6.85), Inches(1.0), Inches(6.2), Inches(6.0),
         COLOR_WHITE, border=COLOR_LIGHT_GRAY)
add_rect(s, Inches(6.85), Inches(1.0), Inches(6.2), Inches(0.4), COLOR_ACCENT)
add_textbox(s, Inches(7.0), Inches(1.04), Inches(6), Inches(0.32),
            "✓ 美國主策略", size=13, bold=True, color=COLOR_WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

us_main = [
    ("來源：[S1×O1]", True, 11),
    ("「Gift-Ready × Japanese Craftsmanship」 — 父親節/聖誕禮盒", True, 11),
    ("", False, 5),
    ("策略描述：", True, 10),
    ("• 主推 B0FL267TCG (Beard Hero) 改為 7-in-1 套組 + 禮盒包裝", False, 9.5),
    ("• Listing A+ 第一屏：『Perfect Gift for Him』", False, 9.5),
    ("• 強調 47 年日系 OEM 工藝故事 (圖文 + 開箱影片)", False, 9.5),
    ("• PPC 預算 60% 集中 Q2 (5-6 月) + Q4 (10-12 月)", False, 9.5),
    ("• 包裝直接可送、不用重新包裝 (差異化 vs Wahl/ANGFAN)", False, 9.5),
    ("", False, 4),
    ("為何優於其他候選：", True, 10),
    ("• 同時對應 ANOVA Top 1 屬性 (gift_suitability_men)", False, 9.5),
    ("• 槓桿 S1+S2+S3 三大優勢", False, 9.5),
    ("• Wahl/ANGFAN 無『日系工藝』故事可說", False, 9.5),
    ("", False, 4),
    ("🔁 支援：[W1×O3] Reddit/Manscaped 替代社群投放", True, 9.5),
    ("⏸ 暫緩：[W5×O1] 母親節女士系列 — 等 Father's Day 跑通", True, 9.5),
    ("❌ 放棄：[W2×T2] Body Groomer OBM — Manscaped 主場放棄", True, 9.5),
    ("", False, 4),
    ("KPI 目標 (12 個月)", True, 10),
    ("• Q2 ROAS 達 4.5+；Beard 類目銷售 +35% YoY", False, 9.5),
]
add_multiline_textbox(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(5.4),
                      us_main, size=9.5, color=COLOR_GRAY)
add_footer(s, 12, TOTAL)


# ============================================================
# Section 05: 日本市場 BMC + SWOT
# ============================================================

# ====== Slide 13: 日本 — 震央 + BMC ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "日本市場 — BMC 震央：產品導向 (Product-Driven)",
               "ANOVA Top 3 屬性皆為產品規格 (附件 F=2630 / 刻度 F=2412 / 電源 F=1827) → 從產品規格回推")

add_rect(s, Inches(0.3), Inches(1.0), Inches(12.7), Inches(0.6), COLOR_ACCENT)
add_textbox(s, Inches(0.45), Inches(1.05), Inches(12.5), Inches(0.5),
            "震央：產品導向 — 日本市場決策深度高，付費關鍵在『刻度精度 / 附件齊全 / 防水 dB 數值』，"
            "從產品 → VP → 顧客/通路鏈最有效",
            size=11, bold=True, color=COLOR_WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_bmc_block(s, x1, margin_y2, col1_w, top_h2, "① KP", [
    ("• 樂天市場 — 台灣廠商稀缺直營資格", False, 8.5),
    ("• Amazon JP (FBA + 鼻毛/鬍鬚類目廣告)", False, 8.5),
    ("• 日本物流商 (1日內到貨能力)", False, 8.5),
    ("• 群眾募資平台 (Makuake / Green Funding)", False, 8.5),
    ("• 跨境官網 urbaner.jp/ja", False, 8.5),
    ("• (待補) 理容師/美容クリニック背書", False, 8.5),
])
add_bmc_block(s, x2, margin_y2, col2_w, half_h2, "② KA", [
    ("• 規格『數值化』Listing 文案", False, 8),
    ("  (附件數/段階/dB/IPX/重量)", False, 8),
    ("• 父の日 + 年末ギフト檔期經營", False, 8),
    ("• 消費者教育文章 (note/価格.com)", False, 8),
])
add_bmc_block(s, x2, margin_y2 + half_h2 + gap, col2_w, half_h2, "③ KR", [
    ("• Hero SKU B0GBWZBMS5 (Nose/Ear)", False, 8),
    ("• 7,163 則 JP 自有評論", False, 8),
    ("• 樂天直營資格 (台廠稀缺)", False, 8),
])
add_bmc_block(s, x3, margin_y2, col3_w, top_h2, "④ VP — 日本版", [
    ("『精度・耐久・分貝数値化』", True, 11),
    ("身嗜み × 自然な仕上がり × 丸洗い", True, 9),
    ("", False, 4),
    ("• アタッチメント 7 個 / 38 段階", False, 8.5),
    ("• 0.5mm 単位調整 / IPX7", False, 8.5),
    ("• 騒音 50dB 以下 (鼻/耳/嬰兒款)", False, 8.5),
    ("• 47 年代工 Panasonic / Philips 經歷", False, 8.5),
    ("• 樂天直営正規品 + 1 年保証", False, 8.5),
    ("• 父の日 のし対応 + 化粧箱", False, 8.5),
    ("• 雙電源策略：乾電池 + 充電並行", False, 8.5),
    ("  (Segment 1 仍偏好乾電池)", False, 8.5),
], color_bar=COLOR_ACCENT, fill=COLOR_BG_LIGHT)
add_bmc_block(s, x4, margin_y2, col4_w, half_h2, "⑤ CR", [
    ("• 樂天 1 年保証承諾", False, 8),
    ("• 跨境官網日文 FAQ", False, 8),
    ("• 知恵袋常見問題追蹤回應", False, 8),
])
add_bmc_block(s, x4, margin_y2 + half_h2 + gap, col4_w, half_h2, "⑥ CH", [
    ("• Amazon JP — amazon.co.jp/URBANER", False, 8),
    ("• 樂天市場 (台廠稀缺直營)", False, 8),
    ("• 跨境官網 + Makuake/Green Funding", False, 8),
])
add_bmc_block(s, x5, margin_y2, col5_w, top_h2, "⑦ CS — 日本", [
    ("Segment 1 (91.6%, n=6,559)", True, 9),
    ("Power Source / 鼻毛 / CP值 / 易用性", False, 8),
    ("avg★ = 3.46  (大眾市場)", False, 8),
    ("→ Listing bullet 1 寫電源類型", False, 8),
    ("", False, 4),
    ("Segment 2 (7.8%, n=561)", True, 9),
    ("Adjustable Comb / 高 spec 玩家", False, 8),
    ("avg★ = 4.01 (高滿意度精品客)", False, 8),
    ("→ Hero SKU 主目標客", False, 8),
    ("", False, 4),
    ("Segment 3 (0.6%, n=43)", True, 9),
    ("Vibration / 馬達 RPM 重度玩家", False, 8),
])
add_bmc_block(s, x1, margin_y2 + top_h2 + gap, Inches(6.31), bot_h2,
              "⑧ 成本結構", [
    ("• Amazon JP / 樂天平台費 + 廣告 (~20%)", False, 8),
    ("• 化粧箱 + のし対応 包裝成本 (檔期)", False, 8),
    ("• 日文在地化 + Note / 価格.com 行銷", False, 8),
    ("• 日本倉儲 (FBA 樂天倉) + 1日內到貨物流", False, 8),
])
add_bmc_block(s, x3 + Inches(0.05), margin_y2 + top_h2 + gap, Inches(6.31), bot_h2,
              "⑨ 收益流", [
    ("• 主：Nose/Ear Trimmer 系列 (Hero B0GBWZBMS5)", False, 8),
    ("• 父の日 (6 月第 3 週日) + 年末ギフト 雙檔期", False, 8),
    ("• Amazon JP + 樂天市場兩通路同步收入", False, 8),
    ("• 替換刃頭 × 耗材包 (decision_depth 高 → recurring 強)", False, 8),
])
add_footer(s, 13, TOTAL)


# ====== Slide 14: 日本 SWOT ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "日本 SWOT — 與 Panasonic / Philips / 日立 / Maxell IZUMI / Royal / Pateker 對比",
               "排除『日本競爭者也有』的項目 (IPX7、丸洗い、附件齊全等共通標配不列入)")

JP_S = [
    ("S1  樂天直営正規品資格 (台廠稀缺)", True, 10),
    ("    Wahl/ANGFAN/Pateker 多無此資格、Panasonic 自有", False, 8.5),
    ("S2  47 年代工日本品牌經歷 (含 Panasonic)", True, 10),
    ("    可說『同樣的工藝、自有品牌』的故事", False, 8.5),
    ("S3  Hero SKU B0GBWZBMS5 已成形", True, 10),
    ("    距理想點 1.97；7,163 則 JP 評論底量", False, 8.5),
    ("S4  雙電源策略 (乾電池 + 充電並行)", True, 10),
    ("    Segment 1 仍偏好乾電池；多數品牌已強制充電化", False, 8.5),
]
JP_W = [
    ("W1  品牌知名度遠低於 Panasonic / 日立 / 皇家", True, 10),
    ("    マイベスト / 価格.com / kakaku 排名缺席", False, 8.5),
    ("W2  缺『理容師監修 / 医療クリニック共同開発』背書", True, 10),
    ("    Panasonic ER-WF41 醫療指定、Rozally 理容師監修", False, 8.5),
    ("W3  電視/家電量販店通路 (ヨドバシ/ビックカメラ) 缺席", True, 10),
    ("    Panasonic 線下展示 + 試用體驗難以對抗", False, 8.5),
    ("W4  Foil Shaver 品類缺席", True, 10),
    ("    Panasonic Lambdash + Braun 9 雙佔", False, 8.5),
    ("W5  日文 SNS / Note / Twitter 經營弱", True, 10),
    ("    口碑擴散通路缺", False, 8.5),
]
JP_O = [
    ("O1  鼻毛電剪市場仍由 IZUMI / 皇家分食 (Maybest #1 IZUMI)", True, 10),
    ("    Panasonic ER-GN71 高價區可由 URBANER 中價切入", False, 8.5),
    ("O2  父の日 (6/15) + 年末ギフト + 中元御中元", True, 10),
    ("    のし対応 是日本特有禮品文化、URBANER 已具備", False, 8.5),
    ("O3  靜音 (38dB→50dB→60dB 競爭) 全品類趨勢", True, 10),
    ("    URBANER 可用 dB 數値化切入", False, 8.5),
    ("O4  日本顧客『専門家監修』偏好強", True, 10),
    ("    與獸醫/理容師合作仍有空間", False, 8.5),
    ("O5  Maxell IZUMI 寡佔鼻毛市場、Panasonic 高價空缺", True, 10),
    ("    URBANER B0GBWZBMS5 中價精品切入", False, 8.5),
]
JP_T = [
    ("T1  Panasonic 全品類絕對王者 + 廣告碾壓", True, 10),
    ("    通路 + 媒體聲量無法對抗", False, 8.5),
    ("T2  日圓持續貶值 (對台幣換匯不利)", True, 10),
    ("    壓縮跨境利潤", False, 8.5),
    ("T3  日本消費者『一機長用』文化 → 替換週期長", True, 10),
    ("    NPI 創新節奏壓力", False, 8.5),
    ("T4  Pateker 寵物電剪日本 No.1、URBANER 難切入", True, 10),
    ("    『20年プロトリマー推薦』口碑深厚", False, 8.5),
    ("T5  價格.com 比較文化 — 規格不夠強會被秒淘汰", True, 10),
]

add_swot_quadrant(s, Inches(0.3), Inches(1.05), quad_w, quad_h,
                  "S — 日本優勢 (vs Panasonic/Philips)", JP_S, COLOR_GREEN)
add_swot_quadrant(s, Inches(6.83), Inches(1.05), quad_w, quad_h,
                  "W — 日本劣勢", JP_W, COLOR_RED)
add_swot_quadrant(s, Inches(0.3), Inches(4.0), quad_w, quad_h,
                  "O — 日本機會", JP_O, COLOR_BLUE)
add_swot_quadrant(s, Inches(6.83), Inches(4.0), quad_w, quad_h,
                  "T — 日本威脅", JP_T, COLOR_AMBER)
add_footer(s, 14, TOTAL)


# ====== Slide 15: 日本 SWOT 策略碰撞 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "日本 SWOT 策略碰撞 — TOWS 矩陣",
               "態勢：S<W (Panasonic 強勢) × O>T (機會仍多) → 主象限：WO 策略聯盟 (補強再進攻)")

jp_so = [
    ("[S1×O2] 樂天直営 × 父の日 のし対応 + 化粧箱 ★ 已可立即執行", True, 9.5),
    ("[S2×O5] 47年代工故事 + Hero SKU 中價精品切入", True, 9.5),
    ("        鼻毛市場 IZUMI #1 與 Panasonic 高價之間", False, 8.5),
    ("[S3×O3] B0GBWZBMS5 加 dB 數値化 (50dB 以下)", True, 9.5),
    ("        靜音切入點", False, 8.5),
    ("[S4×O1] 雙電源策略 — 同款保留乾電池版", True, 9.5),
    ("        守住 Segment 1 (91.6%) 大眾", False, 8.5),
]
jp_st = [
    ("[S1×T2] 樂天直営定價以日圓計、避匯損波動", True, 9.5),
    ("        綁定 1 年保証做價格錨定", False, 8.5),
    ("[S2×T1] 用『Panasonic 同源工藝』故事", True, 9.5),
    ("        把劣勢轉成『同等品質、更合理價格』", False, 8.5),
    ("[S3×T5] Listing 規格表直接對標 ER-GN71、IZN-C240-K", True, 9.5),
    ("        在価格.com 比較文化中正面對決", False, 8.5),
]
jp_wo = [
    ("[W2×O4] 與獸醫/理容師合作做『専門家監修』背書 ★ 主策略候選", True, 9.5),
    ("        補 W2 弱點、抓 O4 文化偏好", False, 8.5),
    ("[W5×O3] 經營 note / Twitter / マイベスト 投稿", True, 9.5),
    ("        圍繞『靜音 38dB』『丸洗い』議題", False, 8.5),
    ("[W1×O2] 父の日 + 年末ギフト 集中投放廣告", True, 9.5),
    ("        以節慶節點滲透市場心智", False, 8.5),
]
jp_wt = [
    ("[W3×T1] 暫不投線下 (ヨドバシ等)、全力線上", True, 9.5),
    ("        避開 Panasonic 通路深度", False, 8.5),
    ("[W4×T1] Foil Shaver 不投自有品牌", True, 9.5),
    ("        改為純 OEM 接單，保留產線", False, 8.5),
    ("[W1×T3] 縮短新品週期、每 12 個月推 1 款 Hero", True, 9.5),
    ("        對抗『一機長用』替換週期", False, 8.5),
]

add_swot_quadrant(s, Inches(0.3), Inches(1.05), quad_w2, quad_h2,
                  "SO — 乘勝追擊", jp_so, COLOR_GREEN)
add_swot_quadrant(s, Inches(6.83), Inches(1.05), quad_w2, quad_h2,
                  "ST — 守株待兔", jp_st, COLOR_AMBER)
add_swot_quadrant(s, Inches(0.3), Inches(4.0), quad_w2, quad_h2,
                  "WO 日本主象限 — 策略聯盟 (補強再進攻)", jp_wo, COLOR_BLUE)
add_swot_quadrant(s, Inches(6.83), Inches(4.0), quad_w2, quad_h2,
                  "WT — 退守/避險", jp_wt, COLOR_RED)
add_footer(s, 15, TOTAL)


# ====== Slide 16: 日本策略定調卡 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "日本 — 策略定調卡",
               "態勢 W>S × O>T → WO 策略聯盟 → 主策略：『専門家監修 × 樂天直営』雙保險")

# 候選評分
add_rect(s, Inches(0.3), Inches(1.0), Inches(6.4), Inches(6.0),
         COLOR_BG_LIGHT, border=COLOR_LIGHT_GRAY)
add_textbox(s, Inches(0.45), Inches(1.1), Inches(6.2), Inches(0.4),
            "日本候選策略評分 (主象限 WO)", size=14, bold=True, color=COLOR_PRIMARY)

jp_cand = [
    ("[W2×O4]  専門家監修 × Hero SKU", True, 10.5),
    ("    可行性 2 / 影響力 3 / 時效性 3 = 8 分 ★ 最高分", False, 9),
    ("    (找 1-2 位理容師/美容師合作即可啟動)", False, 9),
    ("", False, 5),
    ("[W5×O3]  靜音 38dB note/Twitter 社群操作", True, 10.5),
    ("    可行性 3 / 影響力 2 / 時效性 2 = 7 分", False, 9),
    ("", False, 5),
    ("[W1×O2]  父の日 + 年末ギフト 集中投放", True, 10.5),
    ("    可行性 3 / 影響力 2 / 時效性 3 = 8 分 (支援用)", False, 9),
    ("", False, 5),
    ("為何 WO 是日本主象限：", True, 10.5),
    ("• Panasonic/Philips/日立 三強的『品牌+通路+背書』", False, 9),
    ("  優勢碾壓 URBANER 自身能力 (W>S)", False, 9),
    ("• 但鼻毛/送禮文化 + 靜音趨勢 + 樂天直営資格", False, 9),
    ("  仍提供 URBANER 切入空間 (O>T)", False, 9),
    ("• 因此先補強『權威背書 + 社群聲量』(WO)", False, 9),
    ("  再進攻、不可貿然以攻擊性 SO 上場", False, 9),
]
add_multiline_textbox(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.4),
                      jp_cand, size=10, color=COLOR_GRAY)

# 主策略
add_rect(s, Inches(6.85), Inches(1.0), Inches(6.2), Inches(6.0),
         COLOR_WHITE, border=COLOR_LIGHT_GRAY)
add_rect(s, Inches(6.85), Inches(1.0), Inches(6.2), Inches(0.4), COLOR_ACCENT)
add_textbox(s, Inches(7.0), Inches(1.04), Inches(6), Inches(0.32),
            "✓ 日本主策略", size=13, bold=True, color=COLOR_WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

jp_main = [
    ("來源：[W2×O4] + 槓桿 [S1×O2]", True, 11),
    ("「専門家監修 × 樂天直営正規品」雙保險", True, 11),
    ("", False, 4),
    ("策略描述：", True, 10),
    ("• 找 1-2 位日本理容師 / 美容師 / 獸醫合作監修", False, 9.5),
    ("  Hero SKU B0GBWZBMS5 (鼻毛/耳毛) → 補 W2", False, 9.5),
    ("• Listing bullet 1：「専門家監修」徽章", False, 9.5),
    ("• 樂天直営正規品 + 1 年保証 + のし対応", False, 9.5),
    ("• Listing 規格表 (附件 7 個 / 38 段 / 50dB / IPX7)", False, 9.5),
    ("  直接對標 Panasonic ER-GN71 + IZUMI IZN-C240-K", False, 9.5),
    ("• 父の日 (6/15) + 年末ギフト 集中投放", False, 9.5),
    ("• 雙電源版本同步上架 (守 Segment 1 偏好)", False, 9.5),
    ("", False, 4),
    ("為何優於其他候選：", True, 10),
    ("• 同時補 W1 (品牌力)、W2 (背書)、W5 (聲量)", False, 9.5),
    ("• 抓住 O4 (専門家偏好) + O5 (中價空缺)", False, 9.5),
    ("• 槓桿 S1 (樂天直営) + S2 (47 年工藝)", False, 9.5),
    ("", False, 4),
    ("🔁 支援：[W5×O3] note/Twitter 靜音 38dB 內容", True, 9.5),
    ("⏸ 暫緩：[S2×O5] 中價精品 SKU — 等監修上線後", True, 9.5),
    ("❌ 放棄：[W4×T1] Foil Shaver 自有品牌 — Panasonic 主場", True, 9.5),
    ("", False, 4),
    ("KPI 目標 (12 個月)", True, 10),
    ("• 樂天月銷 +50%；価格.com 鼻毛類目進入 Top 10", False, 9.5),
]
add_multiline_textbox(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(5.4),
                      jp_main, size=9.5, color=COLOR_GRAY)
add_footer(s, 16, TOTAL)


# ====== Slide 17: 三市場戰略地圖總覽 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "三市場戰略地圖總覽 — 一頁看懂",
               "整體 / 美國 / 日本 三套主策略並行；BMC 震央與 SWOT 主象限對齊")

# 表頭
headers = [
    (Inches(0.3), Inches(2.4), "層次"),
    (Inches(2.7), Inches(2.4), "BMC 震央"),
    (Inches(5.1), Inches(3.0), "SWOT 主象限與主策略"),
    (Inches(8.1), Inches(2.6), "槓桿 (S/O)"),
    (Inches(10.7), Inches(2.4), "12 個月 KPI 目標"),
]
hdr_y = Inches(1.05)
for x, w, label in headers:
    add_rect(s, x, hdr_y, w, Inches(0.5), COLOR_PRIMARY)
    add_textbox(s, x, hdr_y, w, Inches(0.5), label, size=12, bold=True,
                color=COLOR_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("整體", "資源導向\n(47年OEM\n職人能力)",
     "SO 乘勝追擊\n[S1×O3] 日系工藝×樂天\n×父の日 → 全球品牌敘事",
     "S1 47年OEM\n+ S2 樂天\n+ O3 父の日",
     "Beard+Nose\n雙旗艦\n年營收 +30%"),
    ("美國", "顧客導向\n(送禮買家\nF=1299)",
     "SO 乘勝追擊\n[S1×O1] Father's Day\n禮盒 × 7-in-1 套組",
     "S1 工藝故事\n+ S2 Hero SKU\n+ O1 送禮潮",
     "Q2 ROAS 4.5+\nBeard 類目\n+35% YoY"),
    ("日本", "產品導向\n(附件F=2630\n刻度F=2412)",
     "WO 策略聯盟\n[W2×O4] 專家監修\n× 樂天直営×規格對標",
     "S1 樂天\n+ W2 補背書\n+ O4 専門家",
     "樂天月銷 +50%\n価格.com\n鼻毛 Top 10"),
]
row_h = Inches(1.4)
row_y = Inches(1.6)
row_colors = [COLOR_BG_LIGHT, COLOR_WHITE, COLOR_BG_LIGHT]
for i, row in enumerate(rows):
    y = row_y + row_h * i
    bg = row_colors[i]
    for j, (x, w, _) in enumerate(headers):
        add_rect(s, x, y, w, row_h, bg, border=COLOR_LIGHT_GRAY)
        is_first = (j == 0)
        bold = is_first
        size = 11 if is_first else 9.5
        align = PP_ALIGN.CENTER if is_first else PP_ALIGN.LEFT
        # 合併用 multiline
        lines = [(line, bold, size) for line in row[j].split("\n")]
        add_multiline_textbox(s, x + Inches(0.05), y + Inches(0.08),
                              w - Inches(0.1), row_h - Inches(0.16),
                              lines, size=size, color=COLOR_PRIMARY if is_first else COLOR_GRAY,
                              align=align)

# 底部關鍵原則
add_rect(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(1.0),
         COLOR_BG_LIGHT, border=COLOR_ACCENT, border_w=1.5)
key_lines = [
    ("關鍵原則 (Three Locks)", True, 12),
    ("①  雙引擎不可丟  — OBM + OEM 並行：OBM 是品牌長期、OEM 是現金流與供應鏈彈性", False, 10),
    ("②  Hero SKU 先做深、再做廣  — US: B0FL267TCG / JP: B0GBWZBMS5；不要分散資源", False, 10),
    ("③  避免在競品主場硬碰  — 日本不打 Panasonic Foil Shaver、美國不打 Manscaped Body Groomer", False, 10),
]
add_multiline_textbox(s, Inches(0.5), Inches(6.08), Inches(12.4), Inches(0.9),
                      key_lines, size=10, color=COLOR_PRIMARY)
add_footer(s, 17, TOTAL)


# ====== Slide 18: Now / Next / Later 路線圖 ======
s = prs.slides.add_slide(blank_layout)
add_header_bar(s, "Now / Next / Later — 行動路線圖",
               "0-30 天立即執行 × 31-90 天系統化 × 90 天後策略性投入")

stages = [
    ("Now", "0-30 天", "立即執行 (低成本、高槓桿)", COLOR_RED, [
        ("【整體】整理 47 年代工故事 → 1 支 60 秒品牌片", True, 10),
        ("    + 5 張靜態素材 (US/JP 雙語)", False, 9),
        ("【US】Hero SKU B0FL267TCG Listing A+ 加入", True, 10),
        ("    『Perfect Gift for Him』第一屏", False, 9),
        ("【JP】開始接觸 1-2 位理容師/美容師", True, 10),
        ("    談監修合作，目標 30 天內簽 1 位", False, 9),
        ("【JP】樂天 Listing 加入『正規品 + 1年保証』徽章", True, 10),
        ("【全】評分維度說明：可行性高、可立即啟動", False, 9),
    ]),
    ("Next", "31-90 天", "系統化建置", COLOR_AMBER, [
        ("【US】Father's Day 禮盒打樣 + 包裝定稿", True, 10),
        ("    供應商談判完成，5 月底前到 FBA", False, 9),
        ("【JP】専門家監修文案 + 短影音上線", True, 10),
        ("    Listing bullet 1 同步更新", False, 9),
        ("【JP】開始經營 note + Twitter 帳號", True, 10),
        ("    每週 2 篇『靜音 dB / 丸洗い』內容", False, 9),
        ("【整體】設計『URBANER 直営正規品 + 序號系統』", True, 10),
        ("    對抗山寨平行輸入", False, 9),
    ]),
    ("Later", "90 天 +", "策略性投入", COLOR_GREEN, [
        ("【整體】NPI 週期縮短至 12 個月，", True, 10),
        ("    每年至少 2 款新 Hero SKU", False, 9),
        ("【US】複用素材跑黑五/聖誕 + 開發母親節版", True, 10),
        ("【JP】Hero SKU 雙電源版本同步上架", True, 10),
        ("    守住 Segment 1 (乾電池偏好)", False, 9),
        ("【整體】評估 OEM 越南產線移轉", True, 10),
        ("    對沖美中關稅風險", False, 9),
        ("【整體】Beard/Nose 年營收 > USD 5M 後", True, 10),
        ("    再啟動 Foil Shaver / Body Groomer 自有品牌", False, 9),
    ]),
]

stage_w = Inches(4.2)
stage_h = Inches(5.8)
for i, (name, period, desc, color, items) in enumerate(stages):
    x = Inches(0.3) + (stage_w + Inches(0.15)) * i
    y = Inches(1.05)
    # 主框
    add_rect(s, x, y, stage_w, stage_h, COLOR_WHITE, border=COLOR_LIGHT_GRAY)
    # 上方色條
    add_rect(s, x, y, stage_w, Inches(0.85), color)
    add_textbox(s, x, y + Inches(0.05), stage_w, Inches(0.4),
                name, size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, x, y + Inches(0.45), stage_w, Inches(0.3),
                period, size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, x, y + Inches(0.65), stage_w, Inches(0.25),
                desc, size=9, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # 內容
    add_multiline_textbox(s, x + Inches(0.15), y + Inches(0.95),
                          stage_w - Inches(0.3), stage_h - Inches(1.05),
                          items, size=9.5, color=COLOR_GRAY)
add_footer(s, 18, TOTAL)


# ====== 儲存 ======
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "output", "URBANER_BMC_SWOT.pptx")
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"OK: {out_path}")
print(f"Slides: {len(prs.slides)}")
