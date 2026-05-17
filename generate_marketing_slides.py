"""
Generate 4 marketing question answer slides for URBANER presentation
Answers: 每日雙市場區隔 / 鎖定目標客群 / 競品定位分析 / 消費者偏好分析
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x27, 0x44)
GOLD   = RGBColor(0xC9, 0xA3, 0x6F)
US_BLU = RGBColor(0x2E, 0x5B, 0xFF)
JP_RED = RGBColor(0xD3, 0x2F, 0x4D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF5, 0xF7)
DKGRAY = RGBColor(0x3A, 0x3A, 0x3A)
MID    = RGBColor(0x5A, 0x6A, 0x8A)
GREEN  = RGBColor(0x27, 0xAE, 0x60)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, size=Pt(12), bold=False, color=DKGRAY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "微軟正黑體"
    return txb

def add_label_value(slide, x, y, w, label, value, lsize=Pt(9), vsize=Pt(14),
                    lcolor=MID, vcolor=NAVY):
    add_textbox(slide, x, y, w, Inches(0.3), label, size=lsize, color=lcolor, bold=False)
    add_textbox(slide, x, y+Inches(0.28), w, Inches(0.45), value, size=vsize,
                color=vcolor, bold=True)

def slide_header(slide, title, subtitle="", accent=GOLD):
    # top bar
    add_rect(slide, 0, 0, W, Inches(0.08), fill=accent)
    # title bg
    add_rect(slide, 0, Inches(0.08), W, Inches(0.85), fill=NAVY)
    # title text
    add_textbox(slide, Inches(0.35), Inches(0.12), Inches(10), Inches(0.5),
                title, size=Pt(22), bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(0.58), Inches(10), Inches(0.36),
                    subtitle, size=Pt(11), color=GOLD)
    # slide-number accent line
    add_rect(slide, 0, Inches(0.93), W, Inches(0.025), fill=accent)

def add_card(slide, x, y, w, h, title, lines, title_bg=NAVY, title_color=WHITE,
             body_bg=LGRAY, body_color=DKGRAY, icon=""):
    # card bg
    add_rect(slide, x, y, w, h, fill=body_bg)
    # title strip
    th = Inches(0.44)
    add_rect(slide, x, y, w, th, fill=title_bg)
    add_textbox(slide, x+Inches(0.12), y+Inches(0.06), w-Inches(0.2), th-Inches(0.08),
                (icon+" " if icon else "")+title, size=Pt(11), bold=True, color=title_color)
    # body lines
    lh = Inches(0.28)
    for i, line in enumerate(lines):
        add_textbox(slide, x+Inches(0.14), y+th+Inches(0.06)+i*lh,
                    w-Inches(0.2), lh+Inches(0.04),
                    line, size=Pt(9.5), color=body_color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — 每日雙市場區隔
# ══════════════════════════════════════════════════════════════════════════════
sl1 = prs.slides.add_slide(BLANK)
slide_header(sl1, "每日雙市場區隔",
             "K-Means + PCA 分群分析 | 評論數：US 4,360 則 / JP 7,163 則 | Silhouette US=0.358 / JP=0.570")

# 方法 badge
add_rect(sl1, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.35),
         fill=RGBColor(0xC9, 0xA3, 0x6F))
add_textbox(sl1, Inches(10.55), Inches(0.18), Inches(2.4), Inches(0.28),
            "方法：K-Means (k=3) + PCA", size=Pt(8), color=WHITE, bold=True)

Y0 = Inches(1.05)
CW = Inches(6.0)
GAP= Inches(0.4)
CH = Inches(1.35)

# US column header
add_rect(sl1, Inches(0.28), Y0, CW, Inches(0.42), fill=US_BLU)
add_textbox(sl1, Inches(0.35), Y0+Inches(0.06), CW, Inches(0.3),
            "🇺🇸  美國市場（US）— 3 個消費者群體", size=Pt(13), bold=True, color=WHITE)

# JP column header
add_rect(sl1, Inches(6.85), Y0, CW, Inches(0.42), fill=JP_RED)
add_textbox(sl1, Inches(6.92), Y0+Inches(0.06), CW, Inches(0.3),
            "🇯🇵  日本市場（JP）— 3 個消費者群體", size=Pt(13), bold=True, color=WHITE)

us_segs = [
    ("S1｜禮品購買族", "19.5%  ·  n=850  ·  ★3.68",
     ["核心需求：送禮適合度 × 多功能套組", "代表場景：父親節 / 聖誕節採購",
      "主要類別：Beard / Mustache 87.5%"],
     RGBColor(0xE8, 0xED, 0xFF)),
    ("S2｜大眾主流族（痛點群）", "76.5%  ·  n=3,337  ·  ★3.18",
     ["核心需求：鼻耳毛功能 + 產品耐久", "痛點：滿意度最低，是首要改善對象",
      "主要類別：Beard 94.8%"],
     RGBColor(0xFF, 0xF3, 0xF5)),
    ("S3｜充電式優先族", "4.0%  ·  n=173  ·  ★4.48",
     ["核心需求：充電式電源 × 高端體驗", "滿意度最高，是擴展利基",
      "代表 SKU：B0FL267TCG（Hero SKU）"],
     LGRAY),
]

jp_segs = [
    ("S1｜實用主流族", "91.6%  ·  n=6,559  ·  ★3.46",
     ["核心需求：電源類型＋鼻耳毛＋CP值", "決策重點：充電式 vs 乾電池 是核心分歧",
      "主要類別：Nose/Ear 59.4%"],
     RGBColor(0xFF, 0xF3, 0xF5)),
    ("S2｜精準修剪族", "7.8%  ·  n=561  ·  ★4.01",
     ["核心需求：長度調整精度 × 附件齊全", "滿意度高，是優先攻略的利基客群",
      "代表 SKU：B0GBWZBMS5（Hero SKU）"],
     RGBColor(0xE8, 0xED, 0xFF)),
    ("S3｜低震感需求族", "0.6%  ·  n=43  ·  ★3.05",
     ["核心需求：低震動 × 低噪音", "族群極小、優先度最低",
      "主要類別：Nose/Ear 62.8%"],
     LGRAY),
]

for i, (title, stat, lines, bg) in enumerate(us_segs):
    cy = Y0 + Inches(0.5) + i*(CH + Inches(0.1))
    add_rect(sl1, Inches(0.28), cy, CW, CH, fill=bg)
    # accent left bar
    add_rect(sl1, Inches(0.28), cy, Inches(0.07), CH, fill=US_BLU)
    add_textbox(sl1, Inches(0.42), cy+Inches(0.07), CW-Inches(0.22), Inches(0.28),
                title, size=Pt(11), bold=True, color=US_BLU)
    add_textbox(sl1, Inches(0.42), cy+Inches(0.3), CW-Inches(0.22), Inches(0.22),
                stat, size=Pt(9), color=MID)
    for j, l in enumerate(lines):
        add_textbox(sl1, Inches(0.50), cy+Inches(0.52)+j*Inches(0.24),
                    CW-Inches(0.3), Inches(0.24), "• " + l, size=Pt(9), color=DKGRAY)

for i, (title, stat, lines, bg) in enumerate(jp_segs):
    cy = Y0 + Inches(0.5) + i*(CH + Inches(0.1))
    add_rect(sl1, Inches(6.85), cy, CW, CH, fill=bg)
    add_rect(sl1, Inches(6.85), cy, Inches(0.07), CH, fill=JP_RED)
    add_textbox(sl1, Inches(6.99), cy+Inches(0.07), CW-Inches(0.22), Inches(0.28),
                title, size=Pt(11), bold=True, color=JP_RED)
    add_textbox(sl1, Inches(6.99), cy+Inches(0.3), CW-Inches(0.22), Inches(0.22),
                stat, size=Pt(9), color=MID)
    for j, l in enumerate(lines):
        add_textbox(sl1, Inches(7.07), cy+Inches(0.52)+j*Inches(0.24),
                    CW-Inches(0.3), Inches(0.24), "• " + l, size=Pt(9), color=DKGRAY)

# insight bar bottom
add_rect(sl1, 0, Inches(6.88), W, Inches(0.62), fill=NAVY)
add_textbox(sl1, Inches(0.25), Inches(6.95), W-Inches(0.5), Inches(0.45),
            "關鍵發現：US 最大群體痛點集中在「功能持久度」；JP 市場分群最清晰（silhouette 0.570），"
            "電源偏好是最大決策分歧——兩市場不可用同一行銷策略。",
            size=Pt(10), color=GOLD, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 鎖定目標客群
# ══════════════════════════════════════════════════════════════════════════════
sl2 = prs.slides.add_slide(BLANK)
slide_header(sl2, "鎖定目標客群",
             "ANOVA 優先排序 × 消費者輪廓 × Hero SKU 對應")

Y0 = Inches(1.05)
# 2-col layout: US left, JP right
for col, (market, color, segs) in enumerate([
    ("🇺🇸 美國目標客群", US_BLU, [
        ("Priority Target", "S3 — 充電式優先族",
         ["佔比：4.0% ｜ 滿意度最高 ★4.48", "特徵：注重電源品質、追求高端體驗",
          "行動：Hero SKU B0FL267TCG 主打",
          "場景：日常自用 × 高端男士美容"],
         GREEN),
        ("Secondary Target", "S1 — 禮品購買族",
         ["佔比：19.5% ｜ 平均 ★3.68", "特徵：父親節/聖誕節送禮採購者",
          "行動：Gift Ready 包裝 + 節慶廣告",
          "場景：Q2（父親節）× Q4（聖誕）"],
         GOLD),
        ("Deprioritized", "S2 — 大眾主流族",
         ["佔比：76.5% ｜ 滿意度偏低 ★3.18", "策略：先解決痛點（功能耐久）再擴量",
          "行動：Listing 優化、Review 管理"],
         MID),
    ]),
    ("🇯🇵 日本目標客群", JP_RED, [
        ("Priority Target", "S2 — 精準修剪族",
         ["佔比：7.8% ｜ 滿意度最高 ★4.01", "特徵：要求 0.5mm 精度、附件齊全",
          "行動：Hero SKU B0GBWZBMS5 主打",
          "場景：楽天直営店 × 精品修剪話題"],
         GREEN),
        ("Secondary Target", "S1 — 實用主流族",
         ["佔比：91.6% ｜ 平均 ★3.46", "特徵：CP 值敏感、電源類型為關鍵決策",
          "行動：楽天マラソン × 積分倍數攻勢",
          "場景：日常自用 × 節慶促銷"],
         GOLD),
        ("Deprioritized", "S3 — 低震感族",
         ["佔比：0.6% ｜ 滿意度偏低 ★3.05", "策略：族群過小，暫不主動開拓",
          "行動：自然流量維持即可"],
         MID),
    ]),
]):
    x0 = Inches(0.25) if col == 0 else Inches(6.85)
    cw = Inches(6.0)
    add_rect(sl2, x0, Y0, cw, Inches(0.38), fill=color)
    add_textbox(sl2, x0+Inches(0.12), Y0+Inches(0.06), cw, Inches(0.28),
                market, size=Pt(13), bold=True, color=WHITE)

    for i, (label, seg_name, lines, accent) in enumerate(segs):
        cy = Y0 + Inches(0.46) + i * Inches(1.72)
        # label badge
        add_rect(sl2, x0, cy, cw, Inches(0.30), fill=accent)
        add_textbox(sl2, x0+Inches(0.12), cy+Inches(0.04), Inches(2.5), Inches(0.24),
                    label, size=Pt(9), bold=True, color=WHITE)
        # seg name
        add_rect(sl2, x0, cy+Inches(0.30), cw, Inches(0.38), fill=NAVY)
        add_textbox(sl2, x0+Inches(0.12), cy+Inches(0.34), cw-Inches(0.2), Inches(0.30),
                    seg_name, size=Pt(12), bold=True, color=WHITE)
        # lines
        body_h = Inches(1.04)
        add_rect(sl2, x0, cy+Inches(0.68), cw, body_h, fill=LGRAY)
        for j, l in enumerate(lines):
            add_textbox(sl2, x0+Inches(0.15), cy+Inches(0.72)+j*Inches(0.24),
                        cw-Inches(0.25), Inches(0.24), "• " + l, size=Pt(9), color=DKGRAY)

# bottom bar
add_rect(sl2, 0, Inches(6.88), W, Inches(0.62), fill=NAVY)
add_textbox(sl2, Inches(0.25), Inches(6.95), W-Inches(0.5), Inches(0.45),
            "策略核心：US 以禮品送禮族為成長引擎，JP 以精準修剪族為利基突破口——資源 80% 集中 Hero SKU，20% 維持其他品項。",
            size=Pt(10), color=GOLD, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 競品定位分析
# ══════════════════════════════════════════════════════════════════════════════
sl3 = prs.slides.add_slide(BLANK)
slide_header(sl3, "競品定位分析",
             "MNL 偏好份額模型 | Axis B 品質矩陣 | 88 SKU × 114 屬性")

Y0 = Inches(1.05)

# Market share big numbers
add_rect(sl3, Inches(0.25), Y0, Inches(12.8), Inches(1.25), fill=NAVY)

for pos, (val, label, color) in enumerate([
    ("5.21%",   "URBANER 美國偏好份額", RGBColor(0xC9, 0xA3, 0x6F)),
    ("94.67%",  "競品合計份額（US）",   RGBColor(0xFF, 0x6B, 0x6B)),
    ("1.601",   "US Hero B0FL267TCG 理想點距離", RGBColor(0x4C, 0xC9, 0xF0)),
    ("1.970",   "JP Hero B0GBWZBMS5 理想點距離", RGBColor(0x4C, 0xC9, 0xF0)),
]):
    x = Inches(0.4) + pos * Inches(3.2)
    add_textbox(sl3, x, Y0+Inches(0.1), Inches(3.0), Inches(0.5),
                val, size=Pt(26), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(sl3, x, Y0+Inches(0.62), Inches(3.0), Inches(0.35),
                label, size=Pt(9), color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

# Positioning gap table
add_textbox(sl3, Inches(0.25), Inches(2.42), Inches(6.0), Inches(0.35),
            "URBANER vs 競品 — 關鍵屬性定位對比", size=Pt(12), bold=True, color=NAVY)

rows = [
    ("屬性",            "URBANER 現況",   "競品平均",       "機會"),
    ("IPX7+ 防水",      "✅ 已達標",      "部分競品未達",   "強化 Badge 曝光"),
    ("7 件以上附件",    "✅ 主流規格",    "多數僅 3–5 件",  "件數視覺化攻勢"),
    ("USB-C 充電",      "✅ 配備",        "部分仍用舊接頭", "技術領先強調"),
    ("0.5mm 精度（JP）","✅ B0GBWZBMS5", "競品少見此規格",  "JP 差異化核心"),
    ("禮盒包裝（US）",  "✅ B0FL267TCG", "多數競品無禮盒",  "父親節主攻點"),
]

col_ws = [Inches(2.0), Inches(1.9), Inches(1.9), Inches(1.8)]
col_xs = [Inches(0.25), Inches(2.27), Inches(4.2), Inches(6.12)]
header_bg = NAVY
row_bgs   = [LGRAY, WHITE]*3

for ri, row in enumerate(rows):
    rh = Inches(0.38) if ri == 0 else Inches(0.35)
    ry = Inches(2.82) + (ri-0) * Inches(0.35)
    for ci, (cell, cw, cx) in enumerate(zip(row, col_ws, col_xs)):
        bg = header_bg if ri == 0 else (RGBColor(0xE8, 0xED, 0xFF) if ri % 2 == 1 else WHITE)
        fg = WHITE if ri == 0 else DKGRAY
        add_rect(sl3, cx, ry, cw, rh, fill=bg)
        add_textbox(sl3, cx+Inches(0.06), ry+Inches(0.06), cw-Inches(0.1), rh-Inches(0.08),
                    cell, size=Pt(9), color=fg, bold=(ri==0))

# Right: Top competitor positioning
add_textbox(sl3, Inches(8.2), Inches(2.42), Inches(4.9), Inches(0.35),
            "US 競品理想點距離 Top 5（參考定位）", size=Pt(11), bold=True, color=NAVY)

comp_data = [
    ("B0823PD6XD", "URBANER", "偏好份額 3.88% ｜ ★4.40", US_BLU),
    ("B0FCYBJK6T", "競品 #1",  "偏好份額 3.88% ｜ ★3.37", JP_RED),
    ("B0D97YBP2X", "競品 #2",  "偏好份額 3.88% ｜ ★3.91", MID),
    ("B0BPXVV4T7", "競品 #3",  "偏好份額 3.88% ｜ ★4.88", MID),
    ("B07NW44F6Q", "競品 #4",  "偏好份額 3.88% ｜ ★4.88", MID),
]

for i, (asin, brand, info, col) in enumerate(comp_data):
    cy = Inches(2.82) + i * Inches(0.7)
    add_rect(sl3, Inches(8.2), cy, Inches(0.06), Inches(0.62), fill=col)
    add_rect(sl3, Inches(8.28), cy, Inches(4.8), Inches(0.62),
             fill=RGBColor(0xE8, 0xED, 0xFF) if brand == "URBANER" else LGRAY)
    add_textbox(sl3, Inches(8.35), cy+Inches(0.06), Inches(4.6), Inches(0.24),
                f"{brand}｜{asin}", size=Pt(10), bold=(brand=="URBANER"), color=col)
    add_textbox(sl3, Inches(8.35), cy+Inches(0.30), Inches(4.6), Inches(0.22),
                info, size=Pt(9), color=MID)

# insight
add_rect(sl3, Inches(8.2), Inches(6.32), Inches(4.9), Inches(0.55),
         fill=RGBColor(0xE8, 0xED, 0xFF))
add_textbox(sl3, Inches(8.28), Inches(6.38), Inches(4.7), Inches(0.42),
            "URBANER 在 MNL 選擇集排名第 1！\n偏好份額與強力競品並駕齊驅（3.88%）。",
            size=Pt(9), color=US_BLU, bold=True)

# bottom bar
add_rect(sl3, 0, Inches(6.88), W, Inches(0.62), fill=NAVY)
add_textbox(sl3, Inches(0.25), Inches(6.95), W-Inches(0.5), Inches(0.45),
            "機會缺口：競品普遍缺乏「IPX7 + 7件附件 + USB-C + 禮盒包裝」的完整組合——URBANER 的差異化利基明確。",
            size=Pt(10), color=GOLD, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 消費者偏好分析
# ══════════════════════════════════════════════════════════════════════════════
sl4 = prs.slides.add_slide(BLANK)
slide_header(sl4, "消費者偏好分析",
             "Revealed-Preference Logistic Conjoint | ANOVA 屬性重要度 | 願付溢價（WTP）量化")

Y0 = Inches(1.05)

# US & JP importance bars
for col, (market, color, attrs) in enumerate([
    ("🇺🇸 美國 — 屬性重要度", US_BLU, [
        ("功能合一數（多合一）", 51.5),
        ("充電方式（USB-C）",    10.5),
        ("電源類型",             10.4),
        ("附件件數",             10.4),
        ("防水等級（IPX7+）",    8.1),
        ("價格 CP 值",           4.6),
        ("長度調整段數",         2.7),
        ("機身尺寸",             1.7),
    ]),
    ("🇯🇵 日本 — 屬性重要度", JP_RED, [
        ("價格 CP 值",           16.9),
        ("附件件數",             16.0),
        ("長度調整段數",         15.5),
        ("調整精度（mm）",       14.9),
        ("電源類型",             14.9),
        ("防水等級（IPX7+）",    8.5),
        ("電池續航時間",         8.1),
        ("機身尺寸",             5.2),
    ]),
]):
    x0 = Inches(0.25) if col == 0 else Inches(6.85)
    cw = Inches(6.0)
    add_rect(sl4, x0, Y0, cw, Inches(0.36), fill=color)
    add_textbox(sl4, x0+Inches(0.12), Y0+Inches(0.06), cw, Inches(0.26),
                market, size=Pt(12), bold=True, color=WHITE)

    bar_max = Inches(3.2)
    for i, (attr, imp) in enumerate(attrs):
        cy = Y0 + Inches(0.44) + i * Inches(0.55)
        # attr label
        add_textbox(sl4, x0+Inches(0.08), cy+Inches(0.04), Inches(2.2), Inches(0.28),
                    attr, size=Pt(9), color=DKGRAY)
        # bar bg
        add_rect(sl4, x0+Inches(2.35), cy+Inches(0.07), bar_max, Inches(0.28),
                 fill=RGBColor(0xE8, 0xED, 0xFF) if col==0 else RGBColor(0xFF, 0xEB, 0xED))
        # bar fill
        bar_w = bar_max * (imp / 55.0)
        add_rect(sl4, x0+Inches(2.35), cy+Inches(0.07), bar_w, Inches(0.28), fill=color)
        # pct
        add_textbox(sl4, x0+Inches(2.38)+bar_w, cy+Inches(0.08), Inches(0.6), Inches(0.24),
                    f"{imp}%", size=Pt(9), bold=True, color=color)

# WTP section
add_rect(sl4, 0, Inches(5.58), W, Inches(0.32), fill=GOLD)
add_textbox(sl4, Inches(0.3), Inches(5.63), Inches(6), Inches(0.24),
            "消費者願付溢價（WTP）— Conjoint β 值量化", size=Pt(11), bold=True, color=WHITE)

wtp_data = [
    ("附件件數（7件以上）",   "$79.86",  "¥6,800",  "兩市場共同高 WTP"),
    ("長度調整精度（0.5mm）", "$78.18",  "¥7,961",  "JP 最高 WTP 屬性"),
    ("機身緊湊尺寸",          "$89.81",  "—",        "US 市場獨有高值"),
    ("IPX7 防水",             "顯著",    "¥4,200",  "兩市場必要規格"),
    ("USB-C 充電",            "$65+",    "¥3,500",  "US 第二強驅動"),
]

col_xs2 = [Inches(0.25), Inches(3.8), Inches(5.8), Inches(7.8)]
col_ws2 = [Inches(3.5),  Inches(1.9), Inches(1.9), Inches(5.3)]
hdr_row = ["屬性", "美國 WTP (USD)", "日本 WTP (JPY)", "策略意義"]

for ci, (hdr, cw, cx) in enumerate(zip(hdr_row, col_ws2, col_xs2)):
    add_rect(sl4, cx, Inches(5.93), cw, Inches(0.30), fill=NAVY)
    add_textbox(sl4, cx+Inches(0.06), Inches(5.96), cw, Inches(0.24),
                hdr, size=Pt(9), bold=True, color=WHITE)

for ri, row in enumerate(wtp_data):
    ry = Inches(6.25) + ri * Inches(0.27)
    rbg = LGRAY if ri % 2 == 0 else WHITE
    for ci, (val, cw, cx) in enumerate(zip(row, col_ws2, col_xs2)):
        add_rect(sl4, cx, ry, cw, Inches(0.27), fill=rbg)
        fc = US_BLU if ci == 1 else (JP_RED if ci == 2 else DKGRAY)
        add_textbox(sl4, cx+Inches(0.06), ry+Inches(0.03), cw-Inches(0.08), Inches(0.22),
                    val, size=Pt(8.5), color=fc, bold=(ci in [1, 2]))

# bottom bar
add_rect(sl4, 0, Inches(6.88), W, Inches(0.62), fill=NAVY)
add_textbox(sl4, Inches(0.25), Inches(6.95), W-Inches(0.5), Inches(0.45),
            "核心洞察：US 消費者為「功能合一數」付出最大溢價（重要度 51.5%）；JP 消費者則同等重視「CP值、附件數、精度」三軸——定價與規格決策需分市場差異化。",
            size=Pt(10), color=GOLD, bold=True)

# ── Save ───────────────────────────────────────────────────────────────────────
import os
os.makedirs("output", exist_ok=True)
OUT = "output/URBANER_行銷問題回答.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
