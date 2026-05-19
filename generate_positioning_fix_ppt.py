# -*- coding: utf-8 -*-
"""產生「投影片 12 / 15 SKU 更正版」PPT，供使用者複製貼回原簡報。

更正依據：URBANER 雙市場戰略儀表板（2026-05-05）+ research/own_asins.txt
- US 自家旗艦：B0GL2DKVQH（Beard/Mustache，★3.78，距理想 2.49）
  原投影片誤植 B0FL267TCG（競品 Ufree，美國市場頂尖，★4.52，距理想 1.60）
- JP 自家旗艦：B07CYZH2XC（Beard/Mustache，★3.51，距理想 3.43）
  原投影片誤植 B0GBWZBMS5（Nose/Ear 競品，日本市場頂尖，★4.61，距理想 1.97）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 主題色 ----
BG      = RGBColor(0x1A, 0x1A, 0x1A)
GOLD    = RGBColor(0xF2, 0xB2, 0x33)
WHITE   = RGBColor(0xF0, 0xF0, 0xF0)
GRAY    = RGBColor(0x9A, 0x9A, 0x9A)
RED     = RGBColor(0xE8, 0x50, 0x3A)
T_HEAD  = RGBColor(0x18, 0x6E, 0x6C)   # 表頭深teal
T_BODY  = RGBColor(0x3F, 0xA9, 0xA2)   # 表身teal
T_BODY2 = RGBColor(0x38, 0x9A, 0x94)   # 表身teal(交錯)
FONT    = "Microsoft JhengHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, lines, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """lines: list of (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        # CJK 字型
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
        rPr.append(ea)
    return tb


def add_nav(slide):
    items = ["Context", "Market", "Targeting", "Evidence", "Capability", "Execution"]
    n = len(items)
    left0, top, total_w = 0.55, 0.18, 12.2
    step = total_w / n
    for i, it in enumerate(items):
        cx = left0 + i * step
        if it == "Targeting":
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(cx - 0.05), Inches(top - 0.02),
                                         Inches(1.25), Inches(0.36))
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.fill.background()
            tf = box.text_frame
            tf.margin_top = tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = it
            r.font.name = FONT; r.font.size = Pt(12)
            r.font.bold = True; r.font.color.rgb = BG
            add_text(slide, cx - 0.05, top + 0.34, 1.25, 0.22,
                     [("Positioning", 9, GOLD, False)], align=PP_ALIGN.CENTER)
        else:
            add_text(slide, cx, top, 1.4, 0.3,
                     [(it, 12, GRAY, False)], align=PP_ALIGN.LEFT)


def add_title(slide, num, text):
    add_text(slide, 0.45, 0.62, 12.4, 0.8,
             [(f"{num}｜{text}", 30, GOLD, True)])
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.45), Inches(1.34), Inches(12.45), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = GOLD
    ln.line.fill.background()


def add_column(slide, x, w, header, body_lines):
    add_text(slide, x, 1.55, w, 0.4, [(header, 17, GOLD, True)])
    add_text(slide, x, 2.02, w, 2.1, body_lines, line_spacing=1.12)


def add_table(slide, rows):
    nrows, ncols = len(rows), 3
    left, top, width, height = Inches(0.4), Inches(4.45), Inches(12.5), Inches(2.46)
    gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gf.table
    # 關閉內建樣式
    tbl.first_row = False
    tbl.horz_banding = False
    tbl.columns[0].width = Inches(1.95)
    tbl.columns[1].width = Inches(6.55)
    tbl.columns[2].width = Inches(4.0)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(0.46 if ri == 0 else 0.50)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = T_HEAD
            else:
                cell.fill.fore_color.rgb = T_BODY if ri % 2 == 1 else T_BODY2
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tfc = cell.text_frame
            tfc.word_wrap = True
            p = tfc.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT
            r.font.size = Pt(12 if ri == 0 else 11)
            r.font.bold = (ri == 0 or ci == 0)
            r.font.color.rgb = WHITE
            rPr = r._r.get_or_add_rPr()
            rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))


def add_footer_tag(slide):
    add_text(slide, 11.3, 7.02, 1.7, 0.3,
             [("Positioning", 11, GOLD, True)], align=PP_ALIGN.RIGHT)


# ============ Slide 1：更正說明 ============
s = prs.slides.add_slide(BLANK)
_set_bg(s)
add_text(s, 0.7, 0.7, 12, 0.9, [("本檔用途｜投影片 12 / 15 SKU 更正版", 30, GOLD, True)])
add_text(s, 0.7, 1.7, 12, 5.2, [
    ("原簡報投影片 12、15 把「市場頂尖競品」誤植為 URBANER 自家 Hero SKU。", 16, WHITE, False),
    ("以下第 2、3 頁為更正後版本，可直接複製整頁或文字方塊貼回你的原簡報。", 16, WHITE, False),
    ("", 8, WHITE, False),
    ("● 投影片 12（美國 Positioning）", 17, GOLD, True),
    ("   B0FL267TCG  →  B0GL2DKVQH（美國自家旗艦，Beard/Mustache，★3.78）", 14, WHITE, False),
    ("   距理想點 1.601  →  2.49；「為美國 Hero SKU」→「為美國自家旗艦」", 14, WHITE, False),
    ("   B0FL267TCG 其實是競品 Ufree（美國市場頂尖 ★4.52 / 距理想 1.60），改列為追趕對象", 13, GRAY, False),
    ("", 8, WHITE, False),
    ("● 投影片 15（日本 Positioning）", 17, GOLD, True),
    ("   標題與品類 Nose / Ear  →  Beard / Mustache（自家旗艦實際品類）", 14, WHITE, False),
    ("   B0GBWZBMS5  →  B07CYZH2XC（日本自家旗艦，★3.51）", 14, WHITE, False),
    ("   距理想點 1.97  →  3.43；「為日本 Hero SKU」→「為日本自家旗艦」", 14, WHITE, False),
    ("   B0GBWZBMS5 是 Nose/Ear 競品（日本市場頂尖 ★4.61 / 距理想 1.97），改列為尚未站穩戰場", 13, GRAY, False),
    ("", 8, WHITE, False),
    ("資料來源：URBANER 雙市場戰略儀表板（2026-05-05）、research/own_asins.txt", 12, GRAY, False),
])

# ============ Slide 2：更正後 投影片 12 ============
s = prs.slides.add_slide(BLANK)
_set_bg(s)
add_nav(s)
add_title(s, "24", "美國 Positioning：主打 Beard / Mustache Hero SKU")
add_column(s, 0.45, 4.05, "Hero SKU", [
    ("B0GL2DKVQH", 15, WHITE, True),
    ("品類：Beard / Mustache Trimmers", 12.5, WHITE, False),
    ("定位：Gift-Ready Multi-Function Grooming Kit", 12.5, WHITE, False),
    ("STP 評估：距理想點 2.49，為美國自家旗艦", 12.5, WHITE, False),
    ("市場頂尖：競品 Ufree B0FL267TCG ★4.52 / 距理想 1.60（追趕對象）", 11, GOLD, False),
])
add_column(s, 4.65, 4.0, "推薦主訴", [
    ("Perfect Gift for Him", 12.5, WHITE, False),
    ("5-in-1 / 7-in-1 grooming kit", 12.5, RED, False),
    ("Japanese OEM craftsmanship", 12.5, WHITE, False),
    ("USB-C / IPX7 / charging dock", 12.5, RED, False),
    ("套組完整，送禮不需二次包裝", 12.5, WHITE, False),
])
add_column(s, 8.85, 4.05, "價格定位", [
    ("從 Conair 入門品向上區隔", 12.5, WHITE, False),
    ("建議落在 $60–$120", 12.5, WHITE, False),
    ("定位上高於入門平價品", 12.5, WHITE, False),
    ("但低於 Braun Series 9 Pro 等 $300+ 旗艦級產品", 12.5, WHITE, False),
])
add_table(s, [
    ["Listing 區塊", "建議內容", "對應資料訊號"],
    ["主圖", "右上標明 5–7 合 1／附件件數徽章，產品以禮盒與完整套組呈現",
     "num_grooming_functions F=1193；total_attachments_count F=1015"],
    ["A+Content 第一屏", "放父親節、聖誕節、男士日常修容情境；文案寫 Perfect Gift for Him",
     "gift_suitability_men F=1299"],
    ["規格證明", "USB-C、IPX7、防水全水洗、充電底座、附件對照圖",
     "power_source_type F=941；waterproof F=903"],
    ["競品對標", "避免只對 Conair 打價格，改以 Manscaped／Wahl／Andis 對標功能與形象",
     "BMC/SWOT：美國競爭者基準與排除規則"],
])
add_footer_tag(s)
s.notes_slide.notes_text_frame.text = (
    "更正：B0FL267TCG→B0GL2DKVQH（美國自家旗艦 Beard/Mustache ★3.78 距理想2.49）；"
    "距理想點 1.601→2.49；B0FL267TCG 為競品 Ufree（市場頂尖），改列追趕對象。"
)

# ============ Slide 3：更正後 投影片 15 ============
s = prs.slides.add_slide(BLANK)
_set_bg(s)
add_nav(s)
add_title(s, "27", "日本 Positioning：主打 Beard / Mustache Hero SKU")
add_column(s, 0.45, 4.05, "Hero SKU", [
    ("B07CYZH2XC", 15, WHITE, True),
    ("品類：Beard / Mustache Trimmers", 12.5, WHITE, False),
    ("定位：Precision & Durability Grooming Tool", 12.5, WHITE, False),
    ("STP 評估：距理想點 3.43，為日本自家旗艦", 12.5, WHITE, False),
    ("市場頂尖：B0GBWZBMS5（Nose/Ear 競品）★4.61 / 距理想 1.97 — 尚未站穩戰場", 11, GOLD, False),
])
add_column(s, 4.65, 4.0, "推薦主訴", [
    ("アタッチメント X 個", 12.5, WHITE, False),
    ("長さ調整 X 段階", 12.5, RED, False),
    ("0.5mm / 1mm 刻度", 12.5, WHITE, False),
    ("稼働時間 XX 分", 12.5, WHITE, False),
    ("IPX7 防水 / 騒音 XX dB", 12.5, RED, False),
    ("正規品 + 1 年保証", 12.5, WHITE, False),
])
add_column(s, 8.85, 4.05, "通路與信任", [
    ("樂天市場直営資格是台灣廠商稀缺資源", 12.5, WHITE, False),
    ("可與 Amazon JP、跨境官網形成", 12.5, WHITE, False),
    ("「正規品、保固、日文服務」三層信任", 12.5, WHITE, False),
])
add_table(s, [
    ["Listing 區塊", "建議內容", "對應資料訊號"],
    ["主圖", "以規格表、附件數、段數、刻度、續航、防水等數字為第一層視覺",
     "total_attachments_count F=2630；adjustable_comb_settings F=1827"],
    ["Bullet 1–2", "アタッチメント X 個／長さ調整 X 段階／0.5mm 単位／稼働時間 XX 分",
     "STP Strategy Recommendations 7.2"],
    ["信任證明", "URBANER 直営正規品、1 年保証、のし対応、専門家監修或サロン推奨",
     "BMC/SWOT：樂天直営資格、JP 策略卡"],
    ["產品線", "雙電源並行：保留乾電池款守住 Segment 1，同時推出 USB-C 高規格款",
     "Segment 1 核心屬性為 power_source_type"],
])
add_footer_tag(s)
s.notes_slide.notes_text_frame.text = (
    "更正：標題/品類 Nose/Ear→Beard/Mustache；B0GBWZBMS5→B07CYZH2XC（日本自家旗艦 ★3.51 距理想3.43）；"
    "距理想點 1.97→3.43；B0GBWZBMS5 為 Nose/Ear 競品（市場頂尖），改列尚未站穩戰場。"
)

OUT = "output/URBANER_Positioning_修正_P12_P15.pptx"
prs.save(OUT)
print("saved:", OUT)
